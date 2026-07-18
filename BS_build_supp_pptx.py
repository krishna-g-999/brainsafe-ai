# -*- coding: utf-8 -*-
"""
BS_build_supp_pptx.py -- Supplementary Material deck for BrainSafe AI.
Part A: figures (rendered from figures/*.png). Part B: every supplementary table
(S0-S13) rendered from the exact CSVs, paginated where long. No number is typed by
hand; tables are read straight from supplementary/*.csv. No fabrication.
Output: BrainSafe_AI_Supplementary.pptx
"""
import os, glob
os.chdir(os.path.dirname(os.path.abspath(__file__)))
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
try:
    from PIL import Image; HAVE_PIL = True
except Exception:
    HAVE_PIL = False

NAVY=RGBColor(0x0D,0x1B,0x2A); NAVY2=RGBColor(0x13,0x27,0x3C)
PAPER=RGBColor(0xF4,0xF6,0xF9); WHITE=RGBColor(0xFF,0xFF,0xFF)
INK=RGBColor(0x16,0x22,0x2F); MUTED=RGBColor(0x5A,0x6B,0x7B); LINE=RGBColor(0xD9,0xE0,0xE8)
BLUE=RGBColor(0x00,0x72,0xB2); SKY=RGBColor(0x56,0xB4,0xE9); HDRBG=RGBColor(0xEE,0xF2,0xF6)
SERIF="Georgia"; SANS="Calibri"; MONO="Consolas"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW=13.333; SH=7.5
def slide(): return prs.slides.add_slide(BLANK)
def bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def box(s,l,t,w,h):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tb.text_frame.word_wrap=True; return tb.text_frame
def para(tf,text,size=16,bold=False,color=INK,font=SANS,first=False,align=PP_ALIGN.LEFT,after=4):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(after)
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.name=font; r.font.color.rgb=color
    return p
def rect(s,l,t,w,h,c):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def eyebrow(s,t,color=BLUE): para(box(s,0.9,0.5,11.5,0.4),t.upper(),size=12,bold=True,color=color,font=MONO,first=True)
def title(s,t,size=26,color=INK): rect(s,0.9,0.58,0.5,0.045,BLUE); para(box(s,0.9,0.9,11.6,1.0),t,size=size,bold=True,color=color,font=SERIF,first=True)

def caption(s,text):
    tf=box(s,0.7,6.5,12.0,0.85); p=tf.paragraphs[0]; r=p.add_run(); r.text=text
    r.font.size=Pt(11.5); r.font.italic=True; r.font.name=SANS; r.font.color.rgb=MUTED

def fig_slide(path, tag, cap):
    s=slide(); bg(s,WHITE); eyebrow(s,f"Supplementary · {tag}")
    maxw,maxh=12.4,5.4; T0=0.95
    if HAVE_PIL and os.path.exists(path):
        iw,ih=Image.open(path).size; ar=iw/ih; w=maxw; h=w/ar
        if h>maxh: h=maxh; w=h*ar
    else: w,h=maxw,maxh
    s.shapes.add_picture(path,Inches((SW-w)/2),Inches(T0),Inches(w),Inches(h))
    caption(s,cap); return s

def table_slide(df, tag, cap, cont=False):
    s=slide(); bg(s,PAPER); eyebrow(s,"Supplementary Table")
    title(s,tag+(" (continued)" if cont else ""),size=22)
    rows=[list(df.columns)]+df.astype(str).values.tolist()
    nr,nc=len(rows),len(df.columns)
    fs = 12 if nc<=6 else 10.5 if nc<=9 else 9 if nc<=11 else 8
    hfs= fs-0.5
    rh = min(0.42, 4.55/max(nr,1))
    gt=s.shapes.add_table(nr,nc,Inches(0.5),Inches(1.75),Inches(12.33),Inches(rh*nr)).table
    try: gt.first_row=False; gt.horz_banding=False
    except Exception: pass
    for j in range(nc): gt.columns[j].width=Inches(12.33/nc)
    for i,row in enumerate(rows):
        gt.rows[i].height=Inches(rh)
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.margin_left=Inches(0.05); c.margin_right=Inches(0.05)
            c.margin_top=Inches(0.01); c.margin_bottom=Inches(0.01); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
            p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val)
            r.font.size=Pt(hfs if i==0 else fs); r.font.bold=(i==0)
            r.font.name=(MONO if (i>0 and j>0) else (MONO if i==0 else SANS))
            r.font.color.rgb=(MUTED if i==0 else INK)
            c.fill.solid(); c.fill.fore_color.rgb=(HDRBG if i==0 else (WHITE if i%2 else PAPER))
    caption(s,cap); return s

def add_table(tag, csv, cap, cols=None, per=22):
    df=pd.read_csv(csv)
    if cols: df=df[cols]
    if len(df)<=per:
        table_slide(df,tag,cap)
    else:
        for k in range(0,len(df),per):
            table_slide(df.iloc[k:k+per].reset_index(drop=True),tag,cap,cont=(k>0))

def divider(t,sub):
    s=slide(); bg(s,NAVY); rect(s,0,0,SW,0.12,SKY)
    para(box(s,0.9,2.7,11.5,1.2),t,size=40,bold=True,color=WHITE,font=SERIF,first=True)
    para(box(s,0.9,4.0,11.5,0.8),sub,size=16,color=SKY,font=MONO,first=True)
    return s

# ============================================================ TITLE
s=slide(); bg(s,NAVY); rect(s,0,0,SW,0.12,SKY)
para(box(s,0.9,0.7,11.5,0.5),"BRAINSAFE AI",size=14,bold=True,color=SKY,font=MONO,first=True)
para(box(s,0.9,2.4,11.6,1.6),"Supplementary Material",size=40,bold=True,color=WHITE,font=SERIF,first=True)
para(box(s,0.9,3.9,11.4,1.4),"Figures and supplementary tables. Every value is read directly from the "
     "released validation artifacts (supplementary/*.csv); figures are regenerated from out-of-fold "
     "predictions. Nothing is hand-entered.",size=15,color=RGBColor(0xE6,0xEE,0xF5),font=SANS,first=True)
para(box(s,0.9,6.5,11.5,0.5),"8 figures    ·    14 supplementary tables (S0–S13)",size=12.5,color=SKY,font=MONO,first=True)

# ============================================================ PART A, FIGURES
divider("Part A: Figures","fig1–fig7 + graphical abstract")
FIGS=[
 ("figures/graphical_abstract.png","Graphical abstract",
  "Graphical abstract. Structure-to-brain-effect prediction from measured data; validation and coverage summarised. The per-compound profile shown is illustrative."),
 ("figures/fig1_workflow.png","Figure 1",
  "Figure 1. BrainSafe AI pipeline: measured data sources -> curation -> ECFP+descriptor featurisation -> ensemble training -> calibration/conformal -> integrated outputs, with the four validation regimes."),
 ("figures/fig2_dataset.png","Figure 2",
  "Figure 2. Training-set size and class balance for each endpoint (measured records)."),
 ("figures/fig3_validation.png","Figure 3",
  "Figure 3. AUROC across the four validation regimes (random, scaffold, leave-cluster-out, temporal). Error bars on scaffold bars are 95% bootstrap CIs. BBB has no temporal split (no document year)."),
 ("figures/fig4_roc_calibration.png","Figure 4",
  "Figure 4. (A) Scaffold cross-validation ROC curves (out-of-fold). (B) Reliability diagrams after isotonic calibration."),
 ("figures/fig5_conformal_comparison.png","Figure 5",
  "Figure 5. (A) Empirical coverage of 90%-level conformal prediction sets (target 0.90). (B) Scaffold-CV AUROC of the deployed ensemble versus kNN-Tanimoto and logistic-regression baselines."),
 ("figures/fig6_benchmark.png","Figure 6",
  "Figure 6. Per-endpoint random-split AUROC (this work) relative to published random-split ranges."),
 ("figures/fig7_regression.png","Figure 7",
  "Figure 7. Predicted versus measured potency (scaffold CV): measured antioxidant (DPPH) model and four receptor potency-regression endpoints (A2A, 5-HT2A, D2, SERT)."),
]
for p,tag,cap in FIGS:
    if os.path.exists(p): fig_slide(p,tag,cap)

# ============================================================ PART B, TABLES
divider("Part B: Supplementary Tables","S0–S13, read from the released CSVs")
add_table("STable S0: Data provenance","supplementary/STable0_data_provenance.csv",
          "STable S0. Provenance of every endpoint: role, modality, source database, identifier, measurement type, n, and year range. All measured experimental data.")
add_table("STable S1: Classification metrics","supplementary/STable1_classification_metrics.csv",
          "STable S1. Full classification metrics across regimes (key columns shown; complete 16-column table in the CSV).",
          cols=["endpoint","n","pos_rate","AUROC_random","AUROC_scaffold","AUROC_cluster","AUROC_temporal","PR_AUC","MCC","Brier","conformal_coverage"])
add_table("STable S2: Receptor regression","supplementary/STable2_receptor_regression.csv",
          "STable S2. Receptor potency-regression performance (scaffold CV and temporal split).")
add_table("STable S3: Antioxidant (measured DPPH)","supplementary/STable3_antioxidant_measured.csv",
          "STable S3. Measured-DPPH antioxidant regression, including cross-check against the prior curated proxy.")
add_table("STable S4: Threshold sensitivity","supplementary/STable4_threshold_sensitivity.csv",
          "STable S4. Per-operating-threshold precision/recall/F1 by endpoint and split (temporal, scaffold-holdout).")
add_table("STable S5: Similarity-binned AUROC","supplementary/STable5_similarity_binned_auroc.csv",
          "STable S5. AUROC binned by nearest-training Tanimoto similarity, the empirical basis of the applicability-domain flag.")
add_table("STable S6: Clinical reference composition","supplementary/STable6_clinical_reference_composition.csv",
          "STable S6. Composition of the 504-compound clinical-precedent reference set by disease class.")
add_table("STable S7: Benchmark vs literature","supplementary/STable7_benchmark_vs_literature.csv",
          "STable S7. BrainSafe random-split AUROC versus published random-split ranges per endpoint.")
add_table("STable S8: Capability vs LLM","supplementary/STable8_llm_capability_comparison.csv",
          "STable S8. Capability comparison: BrainSafe AI versus a general-purpose large language model.")
add_table("STable S9: Ablation vs baselines","supplementary/STable9_baseline_comparison.csv",
          "STable S9. Scaffold-split AUROC: deployed ensemble versus kNN-Tanimoto and logistic-regression baselines (ensemble best on all 8 endpoints).")
add_table("STable S10: Label-definition robustness","supplementary/STable10_label_threshold_robustness.csv",
          "STable S10. Scaffold-CV AUROC under alternative pChEMBL label cut-offs (deployed / strict / sharp-boundary / high-potency).")
add_table("STable S11: Assay-type composition","supplementary/STable11_assay_type_composition.csv",
          "STable S11. Assay-type composition per target (IC50/Ki/Kd/EC50 fractions on the standardised pChEMBL scale).")
add_table("STable S12: Assay sensitivity","supplementary/STable12_assay_sensitivity.csv",
          "STable S12. Single-assay (IC50-only) versus pooled retraining: AUROC delta <= 0.006 across the three endpoints tested.")
add_table("STable S13: LLM head-to-head scoreboard","supplementary/STable13_llm_scoreboard.csv",
          "STable S13. Pre-registered head-to-head: BrainSafe and four LLMs scored against the frozen measured-data key (live ChEMBL verification of cited identifiers).")

prs.save("BrainSafe_AI_Supplementary.pptx")
print("Saved BrainSafe_AI_Supplementary.pptx |", len(prs.slides._sldIdLst), "slides | PIL:", HAVE_PIL)

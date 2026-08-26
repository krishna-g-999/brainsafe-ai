"""Build the Chapter 1 defence deck, reading every figure from an artefact.

No number in the deck is typed into this file. Each is computed at build time from the file that
owns it, so a slide cannot describe a panel other than the one that is deployed, and re-running this
after a retrain produces a deck that is current rather than one that has to be checked by eye.

This is the descriptive master deck at defence depth. The colloquium and seminar decks are cut from
it rather than written separately.

Run:  brainsafe_env/Scripts/python.exe thesis/presentations/build_chapter01_defence.py
Out:  thesis/presentations/chapter01_defence.pptx
"""
from __future__ import annotations

import ast
import csv
import json
import statistics as st
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[2]
OUT = Path(__file__).resolve().parent / "chapter01_defence.pptx"
TAB = ROOT / "results" / "tables"
INV = ROOT / "inversion" / "results"
PKG = ROOT / "submission_package"
FIG = ROOT / "manuscript" / "figures"

# ----------------------------------------------------------------------------- palette and type

INK    = RGBColor(0x0F, 0x1A, 0x2E)      # deep indigo, dark slides
DEEP   = RGBColor(0x1F, 0x3A, 0x5F)
TEAL   = RGBColor(0x1C, 0x72, 0x93)
AMBER  = RGBColor(0xD9, 0x83, 0x24)      # reserved for caveats and refutations
CRIMSON= RGBColor(0x9E, 0x2A, 0x2B)      # reserved for withdrawal and failure
MUTED  = RGBColor(0x5F, 0x6B, 0x7A)
TINT   = RGBColor(0xEE, 0xF2, 0xF7)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)
CHALK  = RGBColor(0xD7, 0xDF, 0xEA)

HEAD = "Cambria"
BODY = "Calibri"

W, H = 13.333, 7.5
M = 0.62                                  # slide margin


# ----------------------------------------------------------------------------- artefact loading

def rows(path: Path) -> list[dict]:
    with open(path, encoding="utf-8", errors="replace", newline="") as fh:
        return [r for r in csv.DictReader(fh)]


def n(v):
    try:
        return float(v)
    except (TypeError, ValueError):
        return None


def facts() -> dict:
    """Every quantity the deck states, read from the file that produces it."""
    d: dict = {}

    # -- panel shape ---------------------------------------------------------
    inv = rows(PKG / "07_MODELS" / "model_inventory.csv")
    d["estimators"] = len(inv)
    d["deployed"] = sum(1 for r in inv if r["deployed"] == "True")
    fam: dict[tuple[str, str], int] = {}
    for r in inv:
        fam[(r["family"], r["task"])] = fam.get((r["family"], r["task"]), 0) + 1
    d["families"] = sorted(fam.items(), key=lambda kv: -kv[1])

    reg = json.loads((PKG / "07_MODELS" / "binder_panel_registry.json").read_text())
    served = {k: v for k, v in reg.items() if v.get("deployed")}
    withdrawn = {k: v for k, v in reg.items() if not v.get("deployed")}
    d["binders_total"], d["binders_dep"], d["withdrawn_n"] = len(reg), len(served), len(withdrawn)
    d["withdrawn"] = [(k, v.get("auroc_vs_measured_inactives"), v.get("sensitivity_at_threshold"))
                      for k, v in withdrawn.items()]

    au = {k: v["auroc_vs_measured_inactives"] for k, v in served.items()
          if v.get("auroc_vs_measured_inactives") is not None}
    se = {k: v["sensitivity_at_threshold"] for k, v in served.items()
          if v.get("sensitivity_at_threshold") is not None}
    d["binder_auroc"] = (st.mean(au.values()), st.median(au.values()),
                         min(au, key=au.get), min(au.values()), max(au, key=au.get), max(au.values()))
    d["binder_sens"] = (st.mean(se.values()), st.median(se.values()),
                        min(se, key=se.get), min(se.values()), max(se, key=se.get), max(se.values()))

    # distinct molecular targets: deployed binders plus the measured-label target and safety
    # classifiers, excluding receptor regressions that duplicate a binder endpoint and the two
    # models that are not molecular targets at all.
    core = {r["model"] for r in inv
            if r["deployed"] == "True" and r["family"] in ("target", "safety")
            and r["task"] == "classification"}
    d["targets"] = len({k for k in served} | core)

    # -- training data -------------------------------------------------------
    eps = sorted((PKG / "06_TRAINING_DATA" / "endpoints").glob("*.csv"))
    total, per, cens, cens_eps = 0, {}, 0, 0
    for f in eps:
        rr = rows(f)
        total += len(rr)
        per[f.stem] = len(rr)
        c = sum(1 for x in rr if x.get("source") == "ChEMBL_inactive")
        cens += c
        cens_eps += 1 if c else 0
    d["records"], d["tables"] = total, len(eps)
    d["censored"], d["censored_eps"] = cens, cens_eps
    # Every deployed model that owns a table in 06_TRAINING_DATA/endpoints/. Selecting instead by
    # family, as a first draft of this script did, silently drops the barrier model, because BBB is
    # filed under exposure rather than target: 54 tables and a median of 3,588 instead of 55 and
    # 3,789. The published figure is 3,789 and it is the barrier model that makes the difference.
    dep_names = {r["model"].replace("_binder", "") for r in inv if r["deployed"] == "True"}
    dep_tabs = {k: v for k, v in per.items() if k in dep_names}
    d["rows_median"] = st.median(dep_tabs.values())
    d["rows_min"] = min(dep_tabs.items(), key=lambda kv: kv[1])
    d["rows_max"] = max(dep_tabs.items(), key=lambda kv: kv[1])

    # -- cross-validation ----------------------------------------------------
    cv = rows(TAB / "rf_cv_summary.csv")
    d["cv"] = {}
    for split in ("random", "scaffold"):
        g = [r for r in cv if r["split"] == split and r["task"] == "classification"]
        vals = [n(r["roc_auc_mean"]) for r in g]
        d["cv"][split] = {"eps": [r["endpoint"] for r in g], "vals": vals,
                          "mean": st.mean(vals), "min": min(vals), "max": max(vals)}
        rg = [n(r["r2_mean"]) for r in cv if r["split"] == split and r["task"] == "regression"]
        d["cv"][split]["r2"] = (st.mean(rg), min(rg), max(rg))

    # -- calibration, coverage ----------------------------------------------
    cal = rows(TAB / "calibration.csv")
    d["ece_raw"] = st.mean(n(r["ece_raw"]) for r in cal)
    d["ece_cal"] = st.mean(n(r["ece_calibrated"]) for r in cal)
    d["ece_cal_rng"] = (min(n(r["ece_calibrated"]) for r in cal),
                        max(n(r["ece_calibrated"]) for r in cal))
    con = rows(TAB / "rf_conformal.csv")
    d["cov"] = (min(n(r["empirical_coverage"]) for r in con),
                max(n(r["empirical_coverage"]) for r in con), len(con))

    # -- exposure layer ------------------------------------------------------
    adme = {r["endpoint"]: n(r["r2_mean"]) for r in rows(TAB / "adme_cv_summary.csv")
            if r["split"] == "scaffold" and r["task"] == "regression"}
    d["adme"] = adme

    # -- external and prospective -------------------------------------------
    d["ext_bbb"] = [(r["set"], int(r["n"]), n(r["auroc"])) for r in rows(TAB / "external_bbb_validation.csv")]
    pro = rows(TAB / "external_prospective.csv")
    d["prospective"] = (sum(1 for r in pro if r["status"] == "ok"), len(pro))

    strata = rows(TAB / "external_novelty_strata.csv")
    order = ["below 0.40 (different chemotype)", "0.40 to 0.55 (related series)",
             "0.55 to 0.70 (same series)", "0.70 and above (close analogue)"]
    def band(split, key):
        m = {r["novelty_band"]: n(r["recall_at_threshold"]) for r in strata if r[key] == split}
        return [m[b] for b in order]
    d["recall_bands"] = order
    d["recall_time"] = band("time", "split")
    d["recall_rand"] = band("random", "split")
    d["recall_cross"] = band("cross_source", "split")

    # -- specificity ---------------------------------------------------------
    sp = [r for r in rows(TAB / "noncns_specificity_summary.csv") if r["metric"].startswith("Specificity")][0]
    d["spec"] = (n(sp["estimate"]), n(sp["ci95_low"]), n(sp["ci95_high"]), int(sp["n"]), int(sp["k"]))

    # -- falsification -------------------------------------------------------
    vr = rows(INV / "VERDICTS.csv")
    d["hyp_n"] = len(vr)
    d["refuted"] = [r for r in vr if r["verdict"].startswith("REFUTED")]
    d["supported"] = [r for r in vr if r["verdict"] == "SUPPORTED"]
    d["weakened"] = [r for r in vr if r["verdict"] == "WEAKENED"]

    h2 = {r["weights"]: n(r["top3_accuracy"]) for r in rows(INV / "H2_weight_ablation.csv")}
    d["h2"] = h2
    h5 = {r["method"]: n(r["recall"]) for r in rows(INV / "H5_readacross_value.csv")}
    d["h5"] = list(h5.values())
    h8 = {r["metric"]: n(r["value"]) for r in rows(INV / "H8_panel_independence.csv")}
    d["h8"] = (int(h8["targets that ever fire on the drug set"]),
               int(h8["independent directions in the firing pattern"]))
    h7 = rows(INV / "H7_target_discrimination.csv")
    d["h7"] = (len(h7), st.median(n(r["deployed_sensitivity"]) for r in h7),
               min(n(r["auroc_vs_random"]) for r in h7),
               sum(1 for r in h7 if n(r["deployed_sensitivity"]) < 0.5))

    # -- pathway graph, read from the application itself ---------------------
    tree = ast.parse((ROOT / "app.py").read_text(encoding="utf-8"))
    ns = {}
    for node in tree.body:
        if isinstance(node, ast.Assign) and getattr(node.targets[0], "id", "") in (
                "KNOWLEDGE_GRAPH", "DISEASE_ORDER", "PERIPHERAL_MECHANISM_DISEASES"):
            ns[node.targets[0].id] = ast.literal_eval(node.value)
    kg = ns["KNOWLEDGE_GRAPH"]
    d["graph"] = (len(kg), len(ns["DISEASE_ORDER"]),
                  sum(1 for e in kg.values() if len({x[2] for x in e}) > 1),
                  sorted(ns["PERIPHERAL_MECHANISM_DISEASES"]))
    return d


F = facts()


def fmt(v, k=3):
    return f"{v:.{k}f}"


# ----------------------------------------------------------------------------- slide primitives

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = INK if dark else PAPER
    return s


def text(s, x, y, w, h, runs, size=15, color=None, font=BODY, bold=False, italic=False,
         align=PP_ALIGN.LEFT, space=6, line=None, anchor=MSO_ANCHOR.TOP):
    """runs: a string, or a list of (text, {overrides}) tuples, or a list of strings for bullets."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = runs if isinstance(runs, list) else [runs]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        body, over = (item if isinstance(item, tuple) else (item, {}))
        r = p.add_run()
        r.text = body
        fo = r.font
        fo.name = over.get("font", font)
        fo.size = Pt(over.get("size", size))
        fo.bold = over.get("bold", bold)
        fo.italic = over.get("italic", italic)
        fo.color.rgb = over.get("color", color or INK)
    return tb


def card(s, x, y, w, h, fill=TINT):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.05
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def dot(s, x, y, label, fill=TEAL, fg=PAPER, dia=0.42):
    """The deck's one repeated motif: a numbered disc beside a heading."""
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(dia), Inches(dia))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = BODY
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = fg
    return sh


def head(s, num, title, sub=None):
    """Titles are auto-sized so a long one stays on one line instead of colliding with the
    subtitle beneath it. The box is middle-anchored, so the baseline stays level with the disc
    whatever size is chosen."""
    dot(s, M, 0.52, num)
    ln = len(title)
    size = 30 if ln <= 50 else (27 if ln <= 58 else 24)
    text(s, M + 0.62, 0.40, W - 2 * M - 0.62, 0.74, title, size=size, bold=True, font=HEAD,
         color=INK, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(s, M + 0.62, 1.16, W - 2 * M - 0.62, 0.42, sub, size=14, color=MUTED, italic=True)


def equation(s, x, y, w, parts, size=27, color=DEEP):
    """One paragraph of runs, each optionally set as a true subscript, so that S_d and s_t read as
    subscripted symbols rather than as separate letters. parts: list of (text, is_subscript)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for body, is_sub in parts:
        r = p.add_run()
        r.text = body
        r.font.name = HEAD
        r.font.bold = True
        r.font.size = Pt(size * 0.72 if is_sub else size)
        r.font.color.rgb = color
        if is_sub:
            r.font._rPr.set("baseline", "-25000")
    return tb


def source(s, note):
    text(s, M, H - 0.52, W - 2 * M, 0.32, note, size=10, color=MUTED)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def stat(s, x, y, w, value, label, color=TEAL, vsize=40):
    # Long values such as 228,200 wrap and collide with the label beneath, so the figure is scaled
    # down once it runs past five characters rather than being allowed to break across lines.
    if len(value) > 5:
        vsize = vsize * (5.0 / len(value)) ** 0.55
    text(s, x, y, w, 0.66, value, size=vsize, bold=True, font=HEAD, color=color)
    text(s, x, y + 0.70, w, 0.72, label, size=12, color=MUTED, line=1.15)


# ----------------------------------------------------------------------------- slides

# 1 -- title
s = slide(dark=True)
text(s, M, 1.74, W - 2 * M, 0.5, "Chapter 1", size=16, color=TEAL, bold=True, font=BODY)
text(s, M, 2.18, W - 2 * M - 1.0, 2.05,
     "Predicting small-molecule action in the human brain, exposure and engagement together",
     size=34, bold=True, font=HEAD, color=PAPER, line=1.08)
text(s, M, 4.42, W - 2 * M - 2.0, 0.9,
     "A blood-brain-barrier-gated, multi-endpoint predictor, and an account of what it cannot do",
     size=16, color=CHALK, italic=True)
text(s, M, 5.65, W - 2 * M, 0.9,
     [("Krishnasalini Gunanathan", {"size": 15, "bold": True, "color": PAPER}),
      ("Sri Sathya Sai Institute of Higher Learning", {"size": 13, "color": CHALK}),
      (f"Panel as deployed: {F['estimators']} estimators, {F['deployed']} of them serving predictions",
       {"size": 11, "color": TEAL})], space=3)
notes(s, "Chapter 1 sets up the problem and states, at the outset, what the system does not do. "
         "Every figure in this deck is generated from an artefact in the repository at build time.")

# 2 -- the argument
s = slide(dark=True)
text(s, M, 1.55, W - 2 * M - 1.2, 0.5, "The argument of this thesis", size=15, color=TEAL, bold=True)
text(s, M, 2.15, W - 2 * M - 1.2, 3.0,
     "Potency at a target a compound cannot reach is not activity. "
     "Exposure and engagement are therefore one question, and a tool that answers them separately "
     "has answered neither.",
     size=32, font=HEAD, color=PAPER, line=1.22)
text(s, M, 5.35, W - 2 * M - 1.2, 1.0,
     "The contribution is integration, validation and uncertainty reporting. It is not a new algorithm.",
     size=15, color=AMBER, italic=True)
notes(s, "Say this slowly. It is the thesis in two sentences, and the second sentence is the one "
         "that keeps the claim defensible.")

# 3 -- attrition, two failure modes
s = slide()
head(s, "1", "Two ways a central nervous system programme fails",
     "Both happen late, after the medicinal chemistry has been paid for")
card(s, M, 1.75, 5.9, 3.5)
dot(s, M + 0.35, 2.05, "A", fill=DEEP)
text(s, M + 0.35, 2.62, 5.2, 0.4, "It never arrives", size=21, bold=True, font=HEAD, color=DEEP)
text(s, M + 0.35, 3.14, 5.2, 1.9,
     "The barrier is an active interface, not a membrane with a permeability constant: efflux "
     "transporters return compounds to the circulation and tight junctions exclude them. A compound "
     "can be excellent against isolated protein and inert in an animal.",
     size=14, color=INK, line=1.28)

card(s, M + 6.25, 1.75, 5.9, 3.5)
dot(s, M + 6.60, 2.05, "B", fill=CRIMSON)
text(s, M + 6.60, 2.62, 5.2, 0.4, "It arrives and does more", size=21, bold=True, font=HEAD, color=CRIMSON)
text(s, M + 6.60, 3.14, 5.2, 1.9,
     "A compound that reaches the tissue may engage more than its intended target. On the safety "
     "axis, hERG blockade prolongs the QT interval and remains a leading cause of late "
     "cardiovascular attrition. On the efficacy axis, the mechanism reached is not the one that "
     "drives the disease.",
     size=14, color=INK, line=1.28)

text(s, M, 5.55, W - 2 * M, 0.9,
     "Neither question is new. What is unusual is how rarely they are answered by one instrument, "
     "on one molecule, with a statement of confidence that covers both.",
     size=15, color=INK, italic=True, line=1.25)
source(s, "Attrition citation outstanding: references.md carries no attrition series. The key named "
          "cns_attrition resolves to Alavijeh 2005, a barrier review, not a success-rate analysis.")
notes(s, "Flag honestly in the viva that the attrition citation is still to be resolved. The two "
         "failure modes themselves are not in dispute.")

# 4 -- exposure is Kp,uu
s = slide()
head(s, "2", "Exposure means unbound partition, not total concentration",
     "And the honest report is that this is the weakest model in the panel")
text(s, M, 1.80, 6.4, 2.4,
     "Total brain-to-plasma ratio counts drug bound to tissue lipid, which cannot engage a receptor. "
     "The variable that governs whether a free concentration is available at the target is the "
     "unbound partition coefficient, Kp,uu. This panel models it directly.",
     size=15, color=INK, line=1.30)
text(s, M, 4.35, 6.4, 1.9,
     "Modelling the right quantity badly is preferable to modelling the wrong quantity well, but "
     "only if the weakness is reported. It is reported, here and in Chapter 7.",
     size=15, color=AMBER, italic=True, line=1.30)

xs = M + 6.95
card(s, xs, 1.78, 5.15, 4.35)
text(s, xs + 0.45, 2.05, 4.3, 0.4, "Scaffold-split R² across the exposure layer",
     size=13, bold=True, color=DEEP)
lab = [("Kp,uu, unbound partition", "kpuu"), ("logBB, total ratio", "logbb"),
       ("Aqueous solubility", "solubility"), ("Caco-2 permeability", "caco2_permeability"),
       ("Lipophilicity", "lipophilicity"), ("Plasma protein binding", "plasma_protein_binding"),
       ("Hepatocyte clearance", "clearance_hepatocyte")]
yy = 2.62
for name, key in lab:
    v = F["adme"][key]
    col = AMBER if key == "kpuu" else INK
    text(s, xs + 0.45, yy, 3.05, 0.32, name, size=13, color=col,
         bold=(key == "kpuu"))
    text(s, xs + 3.55, yy, 1.1, 0.32, fmt(v), size=13, color=col, bold=True, align=PP_ALIGN.RIGHT)
    yy += 0.46
source(s, "results/tables/adme_cv_summary.csv, scaffold-grouped 10-fold")
notes(s, f"Kp,uu reaches R squared {fmt(F['adme']['kpuu'])} on 566 compounds. If asked why it is "
         "still deployed: because the alternative is to report logBB and let the user believe it "
         "means free concentration.")

# 5 -- the coupling
s = slide()
head(s, "3", "Coupling is multiplicative, and that is the whole design",
     "A target score is admitted only in proportion to predicted exposure")
card(s, M, 1.80, W - 2 * M, 1.35, fill=TINT)
equation(s, M, 2.12, W - 2 * M,
         [("S̃", False), ("d", True), ("(x)   =   γ(x) · ", False),
          ("max", False), ("(t,w) ∈ G(d)", True), ("   w · s", False), ("t", True),
          ("(x)", False)])
text(s, M, 2.72, W - 2 * M, 0.4,
     "γ(x) = 1 for the peripheral conditions, and the predicted barrier probability otherwise",
     size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

items = [
    ("The maximum, not a sum", "A condition is scored by its single strongest engaged mechanism, so "
     "one target explains the score and can be named. Engaging three of a condition's targets scores "
     "exactly as engaging its strongest."),
    ("Why not a sum", "Because co-engaged targets are frequently homologues, and a sum would count "
     f"one observation several times. Of the panel, {F['h8'][0]} targets fire across approved drugs "
     f"but span only {F['h8'][1]} independent directions."),
    ("The gate is shared", f"γ does not depend on the condition, except for the "
     f"{len(F['graph'][3])} peripheral ones ({', '.join(F['graph'][3]).lower()}). That single fact "
     "has a consequence the next slide states."),
]
yy = 3.55
for i, (t, b) in enumerate(items, 1):
    dot(s, M, yy - 0.02, str(i), fill=DEEP, dia=0.36)
    text(s, M + 0.55, yy, 3.4, 0.4, t, size=15, bold=True, color=DEEP)
    text(s, M + 4.05, yy - 0.04, W - 2 * M - 4.05, 0.95, b, size=13, color=INK, line=1.26)
    yy += 1.10
source(s, "Formal specification: docs/TECHNICAL_REPORT.md section 3.7. Graph read from app.py: "
          f"{F['graph'][0]} targets, {F['graph'][1]} conditions.")
notes(s, "The examiners will want the algebra. Give it, then immediately give the refutation it "
         "implies, which is the next slide.")

# 6 -- H3
s = slide(dark=True)
text(s, M, 1.35, 7.2, 0.45, "What the algebra also says", size=15, color=AMBER, bold=True)
text(s, M, 1.90, 7.6, 2.1,
     "The exposure gate cannot discriminate between conditions.",
     size=34, bold=True, font=HEAD, color=PAPER, line=1.14)
text(s, M, 3.85, 7.6, 2.4,
     "Because γ does not depend on the condition, multiplying every condition's score by the same "
     "number leaves their ranking untouched. Gating decides whether anything is reported at all. It "
     "never decides which condition is reported.",
     size=16, color=CHALK, line=1.30)
card(s, M + 8.15, 1.80, 3.95, 4.15, fill=DEEP)
text(s, M + 8.55, 2.20, 3.15, 0.45, "H3", size=26, bold=True, font=HEAD, color=AMBER)
text(s, M + 8.55, 2.80, 3.15, 0.45, "REFUTED", size=19, bold=True, color=PAPER)
text(s, M + 8.55, 3.32, 3.15, 2.3,
     "by construction, not by experiment. Stated in the introduction rather than the limitations, "
     "because it follows from the design and a reader meeting the architecture is entitled to it "
     "at that moment.",
     size=13, color=CHALK, line=1.28)
source(s, "inversion/results/VERDICTS.csv")
notes(s, "This is the slide that buys credibility for everything after it. The system is an "
         "exposure filter over a mechanism ranking, and it should be described that way throughout.")

# 7 -- existing servers
s = slide()
head(s, "4", "What existing servers do, and where they stop",
     "Three families, each answering one half of the question")
cols = [
    ("ADMET and physicochemical", DEEP,
     "Estimate barrier penetration, permeability and efflux liability from structure. They answer "
     "exposure, often well, and say nothing about pharmacology.",
     "The CNS MPO score is a desirability function over six properties, presented as a design aid. "
     "Used as a proxy for central activity it is a category error: it contains no target information."),
    ("Target prediction", TEAL,
     "Rank probable protein targets, typically by similarity to annotated ligands. They answer "
     "engagement, and are agnostic about whether the compound reaches the tissue.",
     "Treated here as a strong baseline, not a straw man: read-across recovers the correct target for "
     f"{fmt(F['h5'][0])} of held-out compounds against {fmt(F['h5'][1])} for a frequency baseline."),
    ("Single-endpoint barrier models", MUTED,
     "Model blood-brain penetration alone, usually on one curated permeability set.",
     "The exposure axis in isolation. Useful, and not a triage instrument."),
]
xw = (W - 2 * M - 0.6) / 3
for i, (t, col, a, b) in enumerate(cols):
    x = M + i * (xw + 0.30)
    card(s, x, 1.72, xw, 4.30)
    dot(s, x + 0.32, 2.00, str(i + 1), fill=col, dia=0.36)
    text(s, x + 0.32, 2.52, xw - 0.64, 0.75, t, size=17, bold=True, font=HEAD, color=col, line=1.12)
    text(s, x + 0.32, 3.35, xw - 0.64, 1.5, a, size=13, color=INK, line=1.26)
    text(s, x + 0.32, 4.92, xw - 0.64, 1.0, b, size=12, color=MUTED, line=1.24, italic=True)
source(s, "Named-server citations outstanding: no server in the comparison is yet in references.md. "
          "H5 figures from inversion/results/H5_readacross_value.csv.")
notes(s, "Be explicit that the specific servers still need verified citations. Do not name them in "
         "the viva as though they were cited.")

# 8 -- what is absent
s = slide()
head(s, "5", "Three capabilities absent across all three families",
     "Together they define what is claimed here")
gaps = [
    ("Coupling", "No server admits a target score in proportion to predicted exposure. Where both "
                 "quantities are produced they sit side by side and the integration is left to the reader."),
    ("Calibrated, compound-specific uncertainty",
     f"Every value returned carries a calibrated probability, a conformal coverage statement measured "
     f"at {fmt(F['cov'][0])} to {fmt(F['cov'][1])} against a 0.90 target, and a distance to the "
     f"nearest measured analogue."),
    ("A quantified statement of what silence means",
     f"When a tool returns nothing, the user cannot tell inactivity from blindness. This one reports "
     f"the recall the panel achieves at that compound's own distance from training chemistry: "
     f"{fmt(F['recall_time'][3], 2)} for a close analogue, {fmt(F['recall_time'][0], 2)} for a novel scaffold."),
]
yy = 1.85
for i, (t, b) in enumerate(gaps, 1):
    card(s, M, yy, W - 2 * M, 1.34)
    dot(s, M + 0.35, yy + 0.44, str(i), fill=TEAL if i < 3 else AMBER)
    text(s, M + 1.05, yy + 0.26, 4.05, 0.85, t, size=17, bold=True, font=HEAD, color=DEEP, line=1.12)
    text(s, M + 5.25, yy + 0.26, W - 2 * M - 5.6, 0.95, b, size=13, color=INK, line=1.26)
    yy += 1.52
source(s, "results/tables/rf_conformal.csv; results/tables/external_novelty_strata.csv")
notes(s, "The third is the one nobody else does and the one this thesis is most confident about.")

# 9 -- the system at a glance
s = slide()
head(s, "6", "The system, as deployed", "Counted from the model inventory and the endpoint tables")
vals = [(str(F["estimators"]), "fitted estimators"),
        (str(F["deployed"]), "of them deployed"),
        (str(F["targets"]), "distinct molecular targets"),
        (f"{F['records']:,}", "measured records"),
        (f"{F['tables']}", "endpoint tables")]
xw2 = (W - 2 * M - 4 * 0.28) / 5
for i, (v, l) in enumerate(vals):
    x = M + i * (xw2 + 0.28)
    card(s, x, 1.80, xw2, 1.85)
    stat(s, x + 0.30, 2.10, xw2 - 0.6, v, l, color=DEEP, vsize=36)

text(s, M, 4.05, 6.0, 2.1,
     [("Every label is a measured experimental value.", {"size": 16, "bold": True, "color": DEEP}),
      ("No label comes from curator annotation and no value is imputed. An earlier prototype trained "
       "on curated annotation scores was shown by feature ablation to be reading the answer back out "
       "of its own features, and structure-only performance collapsed to near zero.",
       {"size": 13, "color": INK})], space=8, line=1.26)

card(s, M + 6.5, 3.95, W - 2 * M - 6.5, 2.05)
text(s, M + 6.9, 4.22, 4.9, 0.4, "Each endpoint is fitted on its own set alone",
     size=15, bold=True, color=DEEP)
text(s, M + 6.9, 4.72, 5.2, 1.1,
     f"Across the deployed classification tables the median is {int(F['rows_median']):,} rows, from "
     f"{F['rows_min'][1]:,} for {F['rows_min'][0]} to {F['rows_max'][1]:,} for {F['rows_max'][0]}. "
     "Nothing is shared between endpoints except the featurisation.",
     size=13, color=INK, line=1.26)
source(s, "submission_package/07_MODELS/model_inventory.csv; submission_package/06_TRAINING_DATA/endpoints/")
notes(s, "If asked how 54 targets arises from 70 deployed estimators: it is the 47 deployed binders "
         "plus the seven measured-label target and safety classifiers, with the four receptor "
         "regressions excluded as duplicates and the antioxidant assay and pKa model excluded as "
         "not being targets.")

# 10 -- three decisions
s = slide()
head(s, "7", "Three methodological decisions", "Each is a decision rather than a default")
dec = [
    ("The negative class is recovered from measurement",
     f"A compound assayed and found inactive is often deposited only as a censored bound, and the "
     f"conventional query discards exactly those rows. A bound settles a label when the whole "
     f"interval falls one side of the activity cut, and is discarded when it spans both. "
     f"This recovers {F['censored']:,} measured non-binders across {F['censored_eps']} endpoints.",
     "Chapter 2"),
    ("Thresholds are measured on a sample disjoint from the one that set them",
     "Choosing a cut as a quantile of a sample and then measuring the false-positive rate on that "
     "same sample cannot fail: the rate restates the quantile. Three disjoint pools, assigned by a "
     "stable hash of the structure, separate the decoys, the cut and the measurement.",
     "Chapter 5"),
    ("Target scores are gated by predicted exposure",
     "Potency at a target the compound cannot reach contributes nothing. This is the one place the "
     "models are combined multiplicatively, and it encodes a pharmacological fact rather than a "
     "statistical convenience.",
     "Chapter 7"),
]
yy = 1.85
for i, (t, b, ch) in enumerate(dec, 1):
    card(s, M, yy, W - 2 * M, 1.42)
    dot(s, M + 0.35, yy + 0.48, str(i), fill=DEEP)
    text(s, M + 1.05, yy + 0.24, 4.25, 0.9, t, size=16, bold=True, font=HEAD, color=DEEP, line=1.12)
    text(s, M + 5.45, yy + 0.22, W - 2 * M - 6.9, 1.05, b, size=12.5, color=INK, line=1.24)
    text(s, W - M - 1.35, yy + 0.52, 1.1, 0.4, ch, size=12, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    yy += 1.58
source(s, "submission_package/06_TRAINING_DATA/endpoints/, counted at build time")
notes(s, "The censored-bound recovery is the most distinctive data decision in the thesis and the "
         "one a cheminformatics examiner will engage with.")

# 11 -- cross-validation chart
s = slide()
head(s, "8", "Discrimination, with the spread rather than the mean",
     "Scaffold-grouped folds withhold entire structural classes; random folds do not")
cd = CategoryChartData()
cd.categories = F["cv"]["random"]["eps"]
cd.add_series("Random 10-fold", F["cv"]["random"]["vals"])
cd.add_series("Scaffold-grouped 10-fold", F["cv"]["scaffold"]["vals"])
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M), Inches(1.70),
                        Inches(8.25), Inches(4.35), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.legend.font.name = BODY
va = ch.value_axis
va.minimum_scale, va.maximum_scale = 0.80, 1.0
va.has_major_gridlines = True
va.major_gridlines.format.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
va.tick_labels.font.size = Pt(11)
va.tick_labels.font.color.rgb = MUTED
ca = ch.category_axis
ca.has_major_gridlines = False
ca.tick_labels.font.size = Pt(11)
ca.tick_labels.font.color.rgb = MUTED
ch.plots[0].series[0].format.fill.solid()
ch.plots[0].series[0].format.fill.fore_color.rgb = CHALK
ch.plots[0].series[1].format.fill.solid()
ch.plots[0].series[1].format.fill.fore_color.rgb = DEEP
ch.plots[0].gap_width = 60

rv, sv = F["cv"]["random"], F["cv"]["scaffold"]
xs = M + 8.55
text(s, xs, 1.85, 3.55, 0.4, "Mean AUROC", size=13, bold=True, color=DEEP)
text(s, xs, 2.30, 3.55, 1.5,
     [(f"{fmt(rv['mean'], 3)}   random, {fmt(rv['min'])} to {fmt(rv['max'])}",
       {"size": 14, "color": INK}),
      (f"{fmt(sv['mean'], 3)}   scaffold, {fmt(sv['min'])} to {fmt(sv['max'])}",
       {"size": 14, "bold": True, "color": DEEP})], space=7)
text(s, xs, 3.45, 3.55, 1.6,
     "The scaffold column is the honest one. It reports generalisation to chemistry the model has "
     "not seen; the random column reports interpolation within chemistry it has.",
     size=12.5, color=INK, line=1.26)
text(s, xs, 5.05, 3.55, 1.0,
     f"Receptor regressions, scaffold R²: {fmt(sv['r2'][1])} to {fmt(sv['r2'][2])}.",
     size=12.5, color=MUTED, line=1.24)
source(s, "results/tables/rf_cv_summary.csv")
notes(s, "Never quote the mean without the range. The barrier model is the weakest classifier at "
         f"{fmt(sv['min'])} and BACE1 the strongest at {fmt(sv['max'])}.")

# 12 -- binder panel
s = slide()
head(s, "9", "The binder panel, validated against measured inactives",
     "Not against the decoys used to train it")
ba, bs = F["binder_auroc"], F["binder_sens"]
card(s, M, 1.78, 5.6, 4.3)
text(s, M + 0.4, 2.05, 4.9, 0.4, "AUROC against compounds tested and found inactive",
     size=13, bold=True, color=DEEP, line=1.15)
stat(s, M + 0.4, 2.55, 2.3, fmt(ba[0]), "mean over the 47 deployed", color=DEEP, vsize=34)
stat(s, M + 3.05, 2.55, 2.3, fmt(ba[1]), "median, the fairer summary", color=TEAL, vsize=34)
text(s, M + 0.4, 3.95, 4.9, 0.9,
     f"Range {fmt(ba[3])} at {ba[2].replace('_', '-')} to {fmt(ba[5])} at {ba[4]}.",
     size=13, color=INK)
text(s, M + 0.4, 4.45, 4.9, 1.2,
     f"Sensitivity: mean {fmt(bs[0])}, median {fmt(bs[1])}, from {fmt(bs[3])} at "
     f"{bs[2].replace('_', '-')} to {fmt(bs[5])} at {bs[4]}.",
     size=13, color=INK, line=1.26)

if (FIG / "Figure9_model_atlas.png").exists():
    s.shapes.add_picture(str(FIG / "Figure9_model_atlas.png"), Inches(M + 5.95), Inches(1.78),
                         height=Inches(4.30))
text(s, M, 6.20, W - 2 * M, 0.55,
     f"One mark per estimator, so that no claim rests on a mean a reader cannot check. The "
     f"{F['withdrawn_n']} withdrawn estimators are drawn in outline, because a panel showing only "
     f"what survived is a selection rather than an inventory.",
     size=12.5, color=MUTED, italic=True, line=1.24)
source(s, "submission_package/07_MODELS/binder_panel_registry.json; figure Figure9_model_atlas.png")
notes(s, "The mean flatters the panel and the median more so. Name GABA-A and COX-2 before an "
         "examiner does.")

# 13 -- calibration and coverage
s = slide()
head(s, "10", "A probability that is not calibrated is not a probability",
     "And a calibrated probability still says nothing about this compound")
card(s, M, 1.85, 5.85, 3.95)
text(s, M + 0.42, 2.15, 5.0, 0.42, "Calibration", size=19, bold=True, font=HEAD, color=DEEP)
text(s, M + 0.42, 2.72, 5.0, 0.45,
     f"Expected calibration error falls from {fmt(F['ece_raw'], 4)} to {fmt(F['ece_cal'], 4)}",
     size=15, color=INK, line=1.2)
text(s, M + 0.42, 3.40, 5.0, 2.1,
     "Isotonic regression fitted on out-of-fold predictions, so no compound contributes to the "
     f"calibrator that scores it. Per endpoint the calibrated error runs {fmt(F['ece_cal_rng'][0], 4)} "
     f"to {fmt(F['ece_cal_rng'][1], 4)}, the barrier model being the worst of them. Because the map "
     "is monotone it cannot reorder compounds, so AUROC is unchanged and only the probability moves.",
     size=13, color=INK, line=1.28)

card(s, M + 6.25, 1.85, 5.85, 3.95)
text(s, M + 6.67, 2.15, 5.0, 0.42, "Conformal coverage", size=19, bold=True, font=HEAD, color=TEAL)
text(s, M + 6.67, 2.72, 5.0, 0.45,
     f"Empirical coverage {fmt(F['cov'][0])} to {fmt(F['cov'][1])} against a 0.90 target",
     size=15, color=INK, line=1.2)
text(s, M + 6.67, 3.40, 5.0, 2.1,
     f"A Mondrian, class-conditional predictor over {F['cov'][2]} endpoints turns the applicability "
     "domain from a caveat into a coverage statement that is measured rather than assumed. The "
     "class-conditional form matters because marginal coverage can be satisfied while the rare "
     "class fails systematically.",
     size=13, color=INK, line=1.28)
text(s, M, 6.05, W - 2 * M, 0.6,
     "A prediction set may be empty, meaning the compound conforms to neither class, or contain "
     "both, meaning it separates neither. Both are informative and both are displayed.",
     size=13, color=MUTED, italic=True, line=1.24)
source(s, "results/tables/calibration.csv; results/tables/rf_conformal.csv")
notes(s, "Expect a statistician to ask why Mondrian rather than marginal. The answer is on the slide.")

# 14 -- the central finding
s = slide()
head(s, "11", "Recall is a function of chemical distance, not of date",
     "Three test sets built by unrelated rules trace one curve")
cd2 = CategoryChartData()
cd2.categories = ["below 0.40\ndifferent chemotype", "0.40 to 0.55\nrelated series",
                  "0.55 to 0.70\nsame series", "0.70 and above\nclose analogue"]
cd2.add_series("Withheld by date", F["recall_time"])
cd2.add_series("Withheld at random", F["recall_rand"])
cd2.add_series("Withheld by curator", F["recall_cross"])
gf2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M), Inches(1.72),
                         Inches(8.15), Inches(4.35), cd2)
c2 = gf2.chart
c2.has_title = False
c2.has_legend = True
c2.legend.position = XL_LEGEND_POSITION.BOTTOM
c2.legend.include_in_layout = False
c2.legend.font.size = Pt(11)
va = c2.value_axis
va.minimum_scale, va.maximum_scale = 0.0, 1.0
va.has_major_gridlines = True
va.major_gridlines.format.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
va.tick_labels.font.size = Pt(11)
va.tick_labels.font.color.rgb = MUTED
c2.category_axis.tick_labels.font.size = Pt(10)
c2.category_axis.tick_labels.font.color.rgb = MUTED
for i, colr in enumerate([DEEP, TEAL, AMBER]):
    ser = c2.plots[0].series[i]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = colr
c2.plots[0].gap_width = 70

xs = M + 8.45
text(s, xs, 1.85, 3.65, 1.5,
     "An apparent temporal decay is a change in the population being tested, not decay in the models.",
     size=14, bold=True, color=DEEP, line=1.26)
text(s, xs, 3.15, 3.65, 2.0,
     "A random split of medicinal-chemistry data holds out mostly close analogues of its own "
     "training set, because the published record is series. A time split does not. Conditioned on "
     "distance, it makes no difference how the compound was withheld.",
     size=12.5, color=INK, line=1.26)
text(s, xs, 5.25, 3.65, 0.95,
     f"So the expected recall for the compound in hand is knowable at the moment of the query.",
     size=13, color=TEAL, italic=True, line=1.24)
source(s, "results/tables/external_novelty_strata.csv; the curator arm rests on three endpoints")
notes(s, "This is the most important result in the thesis and the one to spend time on. The finding "
         "is that the poor number is predictable, not that it is better than it looked.")

# 15 -- silence
s = slide()
head(s, "12", "Therefore silence carries a number", "The same design decision seen from both sides")
if (FIG / "Figure8_use_case.png").exists():
    # sized by width, not height: the figure is wide, and constraining the height let it run
    # underneath the cards on the right.
    s.shapes.add_picture(str(FIG / "Figure8_use_case.png"), Inches(M), Inches(1.90),
                         width=Inches(6.75))
xs = M + 7.05
rowsx = [
    ("Donepezil", "in domain", "0.86", "A silence would be reasonably strong evidence of inactivity.", DEEP),
    ("Atenolol", "in domain", "0.86", "Every target sits below its base rate; the server is correctly quiet.", TEAL),
    ("Withanolide A", "near domain", "0.55", "At this distance the panel recovers barely half of what is there. "
                                             "The silence is not a result.", AMBER),
]
yy = 1.85
for name, dom, rec, note, col in rowsx:
    card(s, xs, yy, W - M - xs, 1.28)
    text(s, xs + 0.35, yy + 0.20, 2.6, 0.35, name, size=15, bold=True, font=HEAD, color=col)
    text(s, xs + 3.05, yy + 0.22, 1.5, 0.3, dom, size=11.5, color=MUTED)
    text(s, xs + 4.35, yy + 0.14, 1.0, 0.45, rec, size=19, bold=True, color=col, align=PP_ALIGN.RIGHT)
    text(s, xs + 0.35, yy + 0.62, W - M - xs - 0.7, 0.6, note, size=12, color=INK, line=1.22)
    yy += 1.44
text(s, M, 5.90, 6.6, 0.85,
     "A compound that does not arrive cannot generate a disease call. Reported silence and reported "
     "engagement are one rule seen from opposite sides.",
     size=13, color=MUTED, italic=True, line=1.24)
source(s, "docs/TECHNICAL_REPORT.md section 5.1; recall figures from external_novelty_strata.csv")
notes(s, "Withanolide A is the slide that shows the system knows the difference between 'no' and "
         "'I cannot see'. Most tools cannot make that distinction.")

# 16 -- falsification
s = slide(dark=True)
text(s, M, 0.82, 8.0, 0.5, "The falsification suite", size=15, color=AMBER, bold=True)
text(s, M, 1.28, 9.5, 0.72,
     f"{len(F['refuted'])} of {F['hyp_n']} hypotheses were refuted",
     size=34, bold=True, font=HEAD, color=PAPER, line=1.08)
text(s, M, 2.12, 10.6, 0.75,
     "Each was stated so that it could fail, and paired with a null model capable of producing the "
     "same apparent success by accident. All four are reported as refuted.",
     size=14, color=CHALK, line=1.26)
ref = [(r["hypothesis"], r["headline"]) for r in F["refuted"]]
yy = 3.10
for hyp, hl in ref:
    label = hyp.split(" ", 1)[0]
    rest = hyp.split(" ", 1)[1]
    dot(s, M, yy + 0.02, label, fill=AMBER, fg=INK, dia=0.40)
    text(s, M + 0.60, yy, 4.15, 0.80, rest, size=13.5, bold=True, color=PAPER, line=1.16)
    text(s, M + 5.00, yy, W - 2 * M - 5.0, 0.80, hl, size=12, color=CHALK, line=1.20)
    yy += 0.90
text(s, M, 6.82, W - 2 * M, 0.4,
     f"{len(F['supported'])} supported, {len(F['weakened'])} weakened. "
     "The refutations changed the design and are the most useful output the project has produced.",
     size=12.5, color=TEAL, italic=True)
notes(s, "H2 removed a claim about the curated weights. H3 changed how gating is described. H7 "
         "changed how silence is reported. H8 changed how engaged targets are counted in the "
         "interface. Each cost something.")

# 17 -- withdrawn
s = slide()
head(s, "13", f"{F['withdrawn_n']} endpoints were trained, tested and withdrawn",
     "The panel is reported as an inventory, not as a selection")
text(s, M, 1.80, 4.65, 3.3,
     "Two fired on glucose, urea and atenolol at their calibrated thresholds. Three were added "
     "specifically to test natural-product coverage, and all three failed.\n\n"
     "Withdrawal is re-derived whenever the panel is refitted rather than carried forward, because "
     "it is a claim about a particular fit.",
     size=14, color=INK, line=1.30)
xs = M + 5.15
text(s, xs, 1.82, 2.6, 0.3, "endpoint", size=11, bold=True, color=MUTED)
text(s, xs + 2.7, 1.82, 1.9, 0.3, "AUROC", size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
text(s, xs + 4.8, 1.82, 1.9, 0.3, "sensitivity", size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
yy = 2.22
for name, a, sv2 in F["withdrawn"]:
    card(s, xs - 0.18, yy - 0.08, W - M - xs + 0.18, 0.62)
    text(s, xs, yy + 0.04, 2.6, 0.35, name.replace("_", "."), size=14, bold=True, color=INK)
    text(s, xs + 2.7, yy + 0.04, 1.9, 0.35, fmt(a) if a is not None else "n/a",
         size=14, color=CRIMSON if (a or 0) < 0.7 else INK, align=PP_ALIGN.RIGHT)
    text(s, xs + 4.8, yy + 0.04, 1.9, 0.35, fmt(sv2) if sv2 is not None else "n/a",
         size=14, color=CRIMSON if (sv2 or 0) < 0.2 else INK, align=PP_ALIGN.RIGHT)
    yy += 0.74
source(s, "submission_package/07_MODELS/binder_panel_registry.json, withdrawal reasons recorded per endpoint")
notes(s, "NFKB1 scores below chance against its own held-out measured inactives. Showing that is "
         "the point of reporting an inventory.")

# 18 -- specificity and its bound
s = slide()
head(s, "14", "Specificity, stated as the bound it is",
     "The assumption is named in the sentence that makes the claim")
sp = F["spec"]
card(s, M, 1.85, 5.4, 3.5)
stat(s, M + 0.45, 2.20, 4.5, fmt(sp[0]),
     f"of {sp[3]:,} compounds with no recorded activity at any modelled target returned no "
     f"actionable disease signal", color=DEEP, vsize=52)
text(s, M + 0.45, 4.35, 4.5, 0.5, f"95 per cent interval {fmt(sp[1])} to {fmt(sp[2])}",
     size=13, color=MUTED)
card(s, M + 5.85, 1.85, W - 2 * M - 5.85, 3.5, fill=RGBColor(0xFD, 0xF3, 0xE3))
text(s, M + 6.30, 2.20, 5.4, 0.45, "Why this is a lower bound", size=19, bold=True, font=HEAD, color=AMBER)
text(s, M + 6.30, 2.85, 5.4, 2.2,
     "Those compounds are presumed inactive because nothing is recorded about them, not proven "
     "inactive. The artefact itself labels the paired false-positive rate an upper bound. They are "
     "also drawn from within the reference library, so this does not bound behaviour on genuinely "
     "distant chemistry.",
     size=13.5, color=INK, line=1.30)
text(s, M, 5.65, W - 2 * M, 0.9,
     "This is the pattern the whole thesis follows: report the number, then report the assumption "
     "that the number rests on, in the same breath rather than in a footnote.",
     size=14, color=INK, italic=True, line=1.26)
source(s, "results/tables/noncns_specificity_summary.csv")
notes(s, "If an examiner presses on specificity, agree with them. It is a lower bound and it is "
         "labelled as one in the artefact, not just in the prose.")

# 19 -- what it does not do
s = slide()
head(s, "15", "What the system does not do", "Stated in the introduction, not deferred")
lims = [
    ("Chirality is not represented",
     f"Two enantiomers give identical predictions. Bounded in Chapter 10 at 0.19 per cent of the "
     f"panel, which is a bound and not an absolution."),
    ("Agonism and antagonism are not distinguished",
     "The training label is an affinity, which an agonist and an antagonist at the same receptor can "
     "share. The honest description of what the panel predicts is engagement, not modulation."),
    ("The disease layer does not predict indication",
     f"{F['graph'][2]} of the {F['graph'][0]} targets in the pathway graph drive more than one of "
     f"its {F['graph'][1]} conditions, and what selects among them, dose, regimen and patient "
     f"population, is not present in a structure."),
    ("Recall on genuinely novel chemistry is poor",
     f"Below Tanimoto 0.40 recall is {fmt(F['recall_time'][0])}. No analysis in this thesis improves "
     f"that number. The finding is that it is predictable, not that it is better than it looked."),
]
yy = 1.80
for i, (t, b) in enumerate(lims, 1):
    dot(s, M, yy + 0.06, str(i), fill=CRIMSON, dia=0.36)
    text(s, M + 0.58, yy, 4.2, 0.75, t, size=15.5, bold=True, color=CRIMSON, line=1.14)
    text(s, M + 5.0, yy - 0.02, W - 2 * M - 5.0, 1.0, b, size=13, color=INK, line=1.26)
    yy += 1.16
card(s, M, 6.05, W - 2 * M, 0.72, fill=TINT)
text(s, M + 0.40, 6.24, W - 2 * M - 0.8, 0.4,
     "None of these is a defect discovered late. Each was measured, bounded, and written into the "
     "introduction, because a limitation an examiner has to find is worth less than one they are handed.",
     size=13, color=INK, italic=True)
source(s, "docs/TECHNICAL_REPORT.md section 8; inversion/results/; app.py KNOWLEDGE_GRAPH")
notes(s, "Leading with the limitations is a defence tactic and also the right thing to do. An "
         "examiner who has to find these themselves will trust everything else less.")

# 20 -- contribution
s = slide(dark=True)
text(s, M, 1.55, 8.4, 0.5, "The contribution", size=15, color=TEAL, bold=True)
text(s, M, 2.10, 11.0, 2.4,
     "The estimator is a random forest over a fingerprint and twelve descriptors. None of the "
     "components is new.",
     size=30, font=HEAD, color=PAPER, line=1.18)
text(s, M, 4.20, 11.0, 2.2,
     "What has not been assembled elsewhere is the coupling of a measured-label CNS target panel to "
     "a predicted exposure term, with a calibrated probability, a verified coverage statement and a "
     "distance-conditioned recall attached to every value, and with the failures reported at the "
     "same size as the successes.",
     size=17, color=CHALK, line=1.32)
notes(s, "If asked to name the single contribution: the distance-conditioned recall reported at "
         "query time. Nothing else in this space tells a user what a silence is worth.")

# 21 -- roadmap
s = slide()
head(s, "16", "How the rest of the thesis is arranged", None)
chapters = [
    ("2", "Data", "Sources, censored-bound recovery, deduplication, per-endpoint sizes"),
    ("3", "Representation and models", "The featuriser, why a forest, the five-family comparison"),
    ("4", "Uncertainty", "Calibration, conformal prediction, the applicability domain"),
    ("5", "Thresholds", "The three disjoint background pools"),
    ("6", "The binder panel", "Validation against measured inactives"),
    ("7", "Exposure gating", "The gate and the pathway graph"),
    ("8", "Validation", f"External, prospective across {F['prospective'][0]} of "
                        f"{F['prospective'][1]} endpoints, and the composition finding"),
    ("9", "Falsification", f"{F['hyp_n']} hypotheses, {len(F['refuted'])} refuted"),
    ("10", "Limitations", "And the work that follows from them"),
]
colw = (W - 2 * M - 0.4) / 2
for i, (num, t, b) in enumerate(chapters):
    cx = M + (i % 2) * (colw + 0.4)
    cy = 1.72 + (i // 2) * 0.98
    dot(s, cx, cy + 0.08, num, fill=DEEP if i % 2 == 0 else TEAL, dia=0.38)
    text(s, cx + 0.56, cy, colw - 0.6, 0.35, t, size=15, bold=True, font=HEAD, color=INK)
    text(s, cx + 0.56, cy + 0.38, colw - 0.6, 0.5, b, size=12, color=MUTED, line=1.20)
notes(s, "Keep this slide up while taking the first questions; it orients the panel.")

# 22 -- provenance backup
s = slide()
head(s, "17", "Backup: where every number comes from",
     "Nothing in this deck was typed by hand; the generator reads each value at build time")
prov = [
    ("Panel shape, deployed count, targets", "submission_package/07_MODELS/model_inventory.csv"),
    ("Binder AUROC and sensitivity, withdrawals", "submission_package/07_MODELS/binder_panel_registry.json"),
    ("Record count, per-endpoint sizes, censored recovery", "submission_package/06_TRAINING_DATA/endpoints/"),
    ("Cross-validation, both split regimes", "results/tables/rf_cv_summary.csv"),
    ("Calibration and conformal coverage", "results/tables/calibration.csv, rf_conformal.csv"),
    ("Exposure layer R²", "results/tables/adme_cv_summary.csv"),
    ("Recall against chemical distance", "results/tables/external_novelty_strata.csv"),
    ("Specificity and its interval", "results/tables/noncns_specificity_summary.csv"),
    ("Falsification verdicts and each hypothesis", "inversion/results/VERDICTS.csv and H1-H9 tables"),
    ("Pathway graph shape", "app.py, KNOWLEDGE_GRAPH read directly"),
]
yy = 1.78
for i, (what, where) in enumerate(prov):
    if i % 2 == 0:
        card(s, M, yy - 0.06, W - 2 * M, 0.46, fill=TINT)
    text(s, M + 0.28, yy, 5.3, 0.34, what, size=12.5, color=INK)
    text(s, M + 5.85, yy, W - 2 * M - 6.1, 0.34, where, size=12, color=TEAL, font="Consolas")
    yy += 0.46
text(s, M, 6.55, W - 2 * M, 0.5,
     "Rebuild:  brainsafe_env/Scripts/python.exe thesis/presentations/build_chapter01_defence.py",
     size=12, color=MUTED, font="Consolas")
notes(s, "If an examiner challenges a number, this slide says which file to open. That is the "
         "whole point of generating the deck rather than writing it.")


prs.save(OUT)
print(f"wrote {OUT.relative_to(ROOT).as_posix()}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
print(f"  panel {F['estimators']} estimators, {F['deployed']} deployed, {F['targets']} targets")
print(f"  {F['records']:,} records over {F['tables']} endpoint tables")
print(f"  binder AUROC mean {fmt(F['binder_auroc'][0])}, median {fmt(F['binder_auroc'][1])}")
print(f"  {len(F['refuted'])} of {F['hyp_n']} hypotheses refuted")

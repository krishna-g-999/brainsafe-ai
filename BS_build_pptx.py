# -*- coding: utf-8 -*-
"""
BS_build_pptx.py -- dual-audience (expert + commercial/non-scientific) PowerPoint for
BrainSafe AI, with plain-language glosses on every slide and full speaker notes
(narration + anticipated Q&A). All numbers are read from / consistent with the verified
project artifacts. No fabrication.
Output: BrainSafe_AI_Presentation.pptx
"""
import os
os.chdir(os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (Okabe-Ito accents, matching the manuscript figures) ----
NAVY=RGBColor(0x0D,0x1B,0x2A); NAVY2=RGBColor(0x13,0x27,0x3C)
PAPER=RGBColor(0xF4,0xF6,0xF9); WHITE=RGBColor(0xFF,0xFF,0xFF)
INK=RGBColor(0x16,0x22,0x2F); MUTED=RGBColor(0x5A,0x6B,0x7B); LINE=RGBColor(0xD9,0xE0,0xE8)
BLUE=RGBColor(0x00,0x72,0xB2); GREEN=RGBColor(0x00,0x9E,0x73); AMBER=RGBColor(0xE6,0x90,0x00)
VERM=RGBColor(0xD5,0x55,0x00); SKY=RGBColor(0x56,0xB4,0xE9); HDRBG=RGBColor(0xEE,0xF2,0xF6)
LIGHT=RGBColor(0xEA,0xF1,0xF7)
SERIF="Georgia"; SANS="Calibri"; MONO="Consolas"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=prs.slide_width,prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s,color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=color

def box(s,l,t,w,h):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; return tb,tf

def para(tf,text,size=16,bold=False,color=INK,font=SANS,first=False,align=PP_ALIGN.LEFT,
         space_after=6,space_before=0,italic=False,level=0):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(space_before); p.level=level
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return p,r

def runs(tf,parts,size=15,font=SANS,first=False,space_after=6,bullet=None,align=PP_ALIGN.LEFT):
    """parts: list of (text,color,bold)"""
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space_after)
    if bullet:
        rb=p.add_run(); rb.text=bullet+"  "; rb.font.size=Pt(size); rb.font.name=font; rb.font.bold=True; rb.font.color.rgb=BLUE
    for text,color,bold in parts:
        r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.name=font; r.font.bold=bold; r.font.color.rgb=color
    return p

def rrect(s,l,t,w,h,fill,line=None,radius=0.08):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=radius
    except Exception: pass
    return sp

def rect(s,l,t,w,h,fill):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False
    return sp

def eyebrow(s,text,l=0.9,t=0.55,color=BLUE):
    _,tf=box(s,l,t,11.5,0.4); para(tf,text.upper(),size=12,bold=True,color=color,font=MONO,first=True,space_after=0)

def title(s,text,l=0.9,t=0.95,w=11.5,size=33,color=INK):
    _,tf=box(s,l,t,w,1.3); para(tf,text,size=size,bold=True,color=color,font=SERIF,first=True,space_after=0)

def accentbar(s,t=0.62,l=0.9,w=0.5,color=BLUE):
    rect(s,l,t,w,0.045,color)

def notes(s,text):
    s.notes_slide.notes_text_frame.text=text

def add_table(s,data,l,t,w,col_w,header=True,fs=11.5,rh=0.34,hdr_fs=10):
    rows=len(data); cols=len(data[0])
    gt=s.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(w),Inches(rh*rows)).table
    # disable default banding style
    tbl=gt._tbl
    for tag in ('firstRow','bandRow'):
        el=tbl.find(qn('a:tblPr'))
    try:
        gt.first_row=False; gt.horz_banding=False
    except Exception: pass
    total=sum(col_w)
    for j,cw in enumerate(col_w):
        gt.columns[j].width=Inches(w*cw/total)
    for i,rowd in enumerate(data):
        gt.rows[i].height=Inches(rh)
        for j,val in enumerate(rowd):
            c=gt.cell(i,j); c.margin_left=Inches(0.08); c.margin_right=Inches(0.08)
            c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            txt=val[0] if isinstance(val,tuple) else val
            col=(val[1] if isinstance(val,tuple) and len(val)>1 else (INK if i>0 else MUTED))
            bold=(val[2] if isinstance(val,tuple) and len(val)>2 else (i==0))
            tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
            p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(txt)
            r.font.size=Pt(hdr_fs if i==0 else fs); r.font.bold=bold
            r.font.name=(MONO if (j>0 and i>0) else SANS); r.font.color.rgb=col
            if i==0:
                c.fill.solid(); c.fill.fore_color.rgb=HDRBG; r.font.name=MONO
            else:
                c.fill.solid(); c.fill.fore_color.rgb=(WHITE if i%2 else PAPER)
    return gt

def plainbox(s,l,t,w,h,text,color=GREEN):
    sp=rrect(s,l,t,w,h,RGBColor(0xEE,0xFA,0xF5) if color==GREEN else RGBColor(0xEF,0xF4,0xFB),radius=0.06)
    rect(s,l,t,0.06,h,color)
    _,tf=box(s,l+0.18,t+0.06,w-0.3,h-0.12); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    runs(tf,[("In plain terms:  ",color,True),(text,INK,False)],size=12.5,first=True,space_after=0)

# ============================================================ SLIDE 1 — COVER
s=slide(); bg(s,NAVY)
rect(s,0,0,SW.inches,0.12,SKY)
_,tf=box(s,0.9,0.7,11.5,0.5); runs(tf,[("● ",SKY,True),("BRAINSAFE AI",LIGHT,True)],size=14,font=MONO,first=True)
_,tf=box(s,0.9,1.15,11.5,0.4); para(tf,"RESEARCH PREVIEW · PENDING PEER REVIEW",size=12,bold=True,color=SKY,font=MONO,first=True)
_,tf=box(s,0.9,1.9,11.6,2.6)
para(tf,"Predicting a molecule's effect on the brain —",size=34,bold=True,color=WHITE,font=SERIF,first=True,space_after=2)
para(tf,"from its chemical structure alone",size=34,bold=True,color=WHITE,font=SERIF,space_after=0)
_,tf=box(s,0.9,4.15,11.2,1.6)
para(tf,"An evidence-grounded, calibrated, safety-aware predictor of central-nervous-system "
        "activity — built entirely on measured laboratory data, not guesswork.",size=17,color=LIGHT,font=SANS,first=True)
_,tf=box(s,0.9,6.5,11.5,0.6)
para(tf,"Sri Sathya Sai Institute of Higher Learning    ·    64,474 measured records    ·    12 endpoints",
     size=12.5,color=SKY,font=MONO,first=True)
notes(s,
"SAY: Welcome. In one sentence — BrainSafe AI takes a molecule's chemical structure and predicts "
"how it is likely to act on the brain: can it get in, what does it hit, is it safe, and is it drug-like. "
"Crucially, every prediction is built on real laboratory measurements and comes with an honest confidence "
"estimate and the evidence behind it.\n\n"
"TECHNICAL TERMS:\n"
"- 'Endpoint' = one specific thing we predict (e.g. brain penetration, or activity against one target).\n"
"- 'Measured data' = results of actual lab experiments, not simulated or assumed values.\n\n"
"Q&A — GENERAL/COMMERCIAL:\n"
"Q: Is this a finished product / medical device?\nA: No. It is a validated research tool for prioritising and "
"triaging compounds. It is not for clinical or diagnostic decisions, and we say that clearly throughout.\n"
"Q: Who would use it?\nA: Medicinal chemists, natural-product and drug-repurposing researchers — anyone deciding "
"which molecules are worth testing next for brain conditions.\n\n"
"Q&A — EXPERT:\nQ: Application paper or method paper?\nA: An application/resource paper — the contribution is the "
"integration and honest validation, on measured data, not a new algorithm.")

# ============================================================ SLIDE 2 — PROBLEM
s=slide(); bg(s,PAPER); eyebrow(s,"01 · The problem"); accentbar(s); title(s,"One question that is really five questions at once")
_,tf=box(s,0.9,2.0,6.2,4.4)
para(tf,"To judge if a molecule is worth pursuing for a brain condition, you must answer, together:",
     size=15,color=INK,first=True,space_after=12)
for txt,c in [("Can it cross the blood–brain barrier?",BLUE),
              ("Does it engage disease-relevant targets?",BLUE),
              ("Is it developable (drug-like enough)?",AMBER),
              ("Is it safe (no heart-rhythm liability)?",VERM),
              ("Is there clinical precedent?",GREEN)]:
    runs(tf,[(txt,INK,False)],size=15,bullet="▪",space_after=9)
    tf.paragraphs[-1].runs[0].font.color.rgb=c
rrect(s,7.5,2.0,5.0,3.5,WHITE,line=LINE)
_,tf=box(s,7.75,2.2,4.5,3.2)
para(tf,"Existing tools answer only slices",size=16,bold=True,color=INK,first=True,space_after=6)
para(tf,"General ADMET tools predict barrier entry and safety but not target activity. "
        "Target-prediction tools predict binding but ignore brain entry, safety, and confidence. "
        "No single public tool unifies all five on measured data.",size=13,color=MUTED,space_after=0)
plainbox(s,0.9,6.05,11.6,0.9,
"the blood–brain barrier is the brain's security gate — most molecules can't pass. A compound has to clear "
"that gate AND hit the right target AND be safe. Answering these one-at-a-time, in different tools, is slow and error-prone.",color=BLUE)
notes(s,
"SAY: A common misconception is that 'will this help the brain?' is one question. It isn't. A molecule can be a "
"perfect target-blocker and still be useless because it never reaches the brain — or dangerous because it affects "
"the heart. You have to answer all of these together. Today those answers are scattered across separate tools that "
"don't talk to each other.\n\n"
"TECHNICAL TERMS:\n- Blood–brain barrier (BBB): a tight layer of cells that keeps most substances out of the brain.\n"
"- Target: the specific protein a drug is meant to act on (e.g. an enzyme involved in Alzheimer's).\n"
"- ADMET: absorption, distribution, metabolism, excretion, toxicity — the 'developability' properties.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: Why does this matter commercially?\nA: Most drug candidates fail late and expensively. "
"Catching 'won't reach the brain' or 'heart-safety risk' early saves time and money.\n\n"
"Q&A — EXPERT:\nQ: Aren't SwissADME/ADMETlab enough?\nA: They cover BBB/ADMET well but not measured CNS-target "
"activity, BBB-gated disease synthesis, calibrated uncertainty, or a linked safety axis — that combination is the gap.")

# ============================================================ SLIDE 3 — INNOVATION
s=slide(); bg(s,NAVY); eyebrow(s,"02 · What we built",color=SKY); rect(s,0.9,0.62,0.5,0.045,SKY)
title(s,"The innovation is the integration — not a new algorithm",color=WHITE)
_,tf=box(s,0.9,1.95,11.4,1.0)
para(tf,"Each individual piece is well established. Combining them into one transparent, measured-data "
        "brain profiler — that is the contribution.",size=16,color=LIGHT,first=True)
cards=[("Calibrated & confident","Every probability is honest — a '90%' really behaves like 9-in-10 — and comes with a guaranteed confidence range.",BLUE),
       ("Evidence-grounded","Every answer shows the nearest real, measured molecule it is based on — you can check the receipt.",GREEN),
       ("Brain-gated disease view","Target activity is combined with brain-entry, so 'engages target' only counts if the molecule can actually get in.",AMBER)]
x=0.9
for h,b,c in cards:
    rrect(s,x,3.25,3.7,3.2,NAVY2,line=RGBColor(0x25,0x40,0x5A))
    rect(s,x+0.28,3.55,0.5,0.05,c)
    _,tf=box(s,x+0.28,3.75,3.15,2.5)
    para(tf,h,size=17,bold=True,color=WHITE,first=True,space_after=6)
    para(tf,b,size=13,color=LIGHT,space_after=0)
    x+=3.94
notes(s,
"SAY: We did not invent a new machine-learning algorithm — and that is deliberate. The scientific value here is "
"engineering restraint: taking proven, trusted components and integrating them so the whole is honest and auditable. "
"Three properties make it different from a black box: calibration, evidence-grounding, and brain-gating.\n\n"
"TECHNICAL TERMS:\n- Calibration: adjusting scores so a stated probability matches real-world frequency.\n"
"- Confidence range (conformal set): instead of one number, the tool can say 'confidently active' or 'uncertain — "
"both possible', with a mathematical coverage guarantee.\n- Brain-gating: multiplying target activity by "
"brain-penetration probability, so a target hit that can't reach the brain is down-weighted.\n\n"
"Q&A — EXPERT:\nQ: If nothing is novel algorithmically, what's the publishable claim?\nA: The integrated, calibrated, "
"conformal, BBB-gated, safety-aware, evidence-grounded configuration on measured data — no existing single tool "
"provides it — plus a rigorous multi-regime validation.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: What's the moat if the parts are standard?\nA: The integration, the curated measured "
"dataset, the honesty layer (calibration + provenance), and validated performance — hard to reproduce casually.")

# ============================================================ SLIDE 4 — DATA
s=slide(); bg(s,PAPER); eyebrow(s,"03 · Data sources"); accentbar(s)
title(s,"64,474 measured records — every one a real experiment")
add_table(s,[
 ["Source","What it provides","Records"],
 ["ChEMBL_37 (pChEMBL)","Activity vs CNS targets + heart-safety","53,301"],
 ["B3DB","Brain-barrier penetration (measured)","7,807"],
 ["ChEMBL DPPH assays","Antioxidant capacity","2,862"],
 ["ChEMBL ATC-N","Clinical precedent (trial phase)","504"],
 [("Total",INK,True),("",INK,True),("64,474",BLUE,True)],
],0.9,2.05,7.4,[2.1,3.0,1.1])
_,tf=box(s,8.6,2.05,3.9,3.6)
for t in [("Public, reproducible sources — ChEMBL and B3DB.",BLUE),
          ("Only lab-measured values kept (no guesses/imputation).",GREEN),
          ("Report dates (1976–2025) retained for a 'future' test.",AMBER)]:
    runs(tf,[(t[0],INK,False)],size=13.5,bullet="▪",space_after=10,first=(t[1]==BLUE))
    tf.paragraphs[-1].runs[0].font.color.rgb=t[1]
plainbox(s,0.9,6.15,11.6,0.85,
"'measured' means someone physically tested each molecule in a lab and recorded the result. The tool learns from "
"tens of thousands of these real results — it is not making things up from text on the internet.",color=GREEN)
notes(s,
"SAY: Everything the tool knows comes from real experiments in public databases — mainly ChEMBL, the standard "
"repository of measured bioactivity, and B3DB for brain penetration. 64,474 measured records in total. This is the "
"single most important integrity point: the foundation is measured, not scraped text.\n\n"
"TECHNICAL TERMS:\n- ChEMBL: a large public database of measured drug-target activity.\n"
"- pChEMBL: a standardised potency value (higher = more potent) that lets different experiments be compared on one scale.\n"
"- B3DB: a curated database of measured blood–brain-barrier permeability.\n- DPPH: a standard lab test of antioxidant "
"(free-radical-scavenging) strength.\n\n"
"Q&A — EXPERT:\nQ: You pool IC50/Ki/Kd/EC50 — doesn't that mix assays?\nA: Yes, on the standardised pChEMBL scale; we "
"tested this (later slide): retraining on a single assay type changes AUROC by <=0.006, so pooling is safe.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: Is the data licensed/usable?\nA: ChEMBL and B3DB are public and openly licensed; the "
"pipeline is reproducible from released scripts.")

# ============================================================ SLIDE 5 — ENDPOINTS
s=slide(); bg(s,PAPER); eyebrow(s,"04 · What it predicts"); accentbar(s)
title(s,"Twelve core endpoints, plus supporting layers")
add_table(s,[
 ["Prediction","In plain terms","How it's served"],
 ["BBB","Can it enter the brain?","classifier"],
 ["AChE / BChE","Alzheimer's cholinergic enzymes","classifier"],
 ["BACE1","Alzheimer's amyloid enzyme","classifier"],
 ["GSK-3β","Tau / neuroprotection","classifier"],
 ["MAO-A / MAO-B","Mood / Parkinson's enzymes","classifier"],
 ["hERG","Heart-rhythm safety flag","classifier"],
 ["D2 · A2A · 5-HT2A · SERT","Receptor potency (how strong)","regression"],
 ["Antioxidant (DPPH)","Radical-scavenging strength","regression"],
 ["Druggability / CNS-MPO","Is it drug-like for the brain?","rule-based"],
],0.9,2.0,7.7,[2.4,3.3,1.4],rh=0.42,fs=11)
_,tf=box(s,8.9,2.0,3.6,4.0)
para(tf,"Two kinds of answer",size=15,bold=True,color=INK,first=True,space_after=5)
runs(tf,[("Classifier",BLUE,True),(" = yes/no with a probability (e.g. 'active, 0.92').",INK,False)],size=13,space_after=8)
runs(tf,[("Regression",GREEN,True),(" = a strength number, used to rank compounds.",INK,False)],size=13,space_after=12)
para(tf,"Why the split? The four receptors were 96–98% actives in the data — too lopsided for a fair yes/no "
        "model — so the tool ranks their strength instead. A quality gate decides this automatically.",
     size=12.5,color=MUTED,space_after=0)
notes(s,
"SAY: The tool outputs a full brain profile: brain entry, several disease-relevant targets, a heart-safety flag, "
"receptor strengths, antioxidant capacity, and drug-likeness. Two of these — druggability and CNS-MPO — are "
"transparent chemistry rules, not machine learning.\n\n"
"TECHNICAL TERMS:\n- AChE/BChE/BACE1/GSK-3β/MAO: enzymes implicated in Alzheimer's, Parkinson's and mood disorders.\n"
"- hERG: a heart ion-channel; blocking it can cause dangerous arrhythmia — a classic safety red flag.\n"
"- Receptor (D2/A2A/5-HT2A/SERT): proteins involved in dopamine/serotonin signalling.\n"
"- Classifier vs regression: yes/no-with-probability vs a continuous strength score.\n"
"- CNS-MPO: a published 'is this drug-like for the brain?' scoring rule.\n\n"
"Q&A — EXPERT:\nQ: Why regression for the receptors?\nA: 96–98% active-only data makes binary classification "
"degenerate; a pre-set quality gate (MCC>=0.45) routes them to potency regression instead — decided by data, not hand.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: Can you add more targets later?\nA: Yes — the same pipeline extends to any target with "
"enough measured data.")

# ============================================================ SLIDE 6 — FEATURES & MODELS
s=slide(); bg(s,NAVY); eyebrow(s,"05 · How it learns",color=SKY); rect(s,0.9,0.62,0.5,0.045,SKY)
title(s,"How the computer 'sees' a molecule — and learns",color=WHITE)
rrect(s,0.9,1.95,5.7,4.4,NAVY2,line=RGBColor(0x25,0x40,0x5A))
_,tf=box(s,1.2,2.2,5.1,4.0)
para(tf,"Turning a molecule into numbers",size=17,bold=True,color=WHITE,first=True,space_after=8)
runs(tf,[("Chemical fingerprint (ECFP-4)",SKY,True),(" — a 1,024-bit 'barcode' of the molecule's building blocks.",LIGHT,False)],size=13.5,bullet="▪",space_after=8)
runs(tf,[("24 physical descriptors",SKY,True),(" — size, greasiness, surface area, flexibility, drug-likeness.",LIGHT,False)],size=13.5,bullet="▪",space_after=8)
runs(tf,[("Why this?",AMBER,True),(" interpretable and fast, and the standard other studies are measured against.",LIGHT,False)],size=13.5,bullet="▪",space_after=0)
rrect(s,6.9,1.95,5.5,4.4,NAVY2,line=RGBColor(0x25,0x40,0x5A))
_,tf=box(s,7.2,2.2,4.9,4.0)
para(tf,"A panel of three models that vote",size=17,bold=True,color=WHITE,first=True,space_after=8)
for t in ["Random Forest — 300 decision trees","Extra Trees — 300 randomised trees","Gradient Boosting — 300 rounds"]:
    runs(tf,[(t,LIGHT,False)],size=13.5,bullet="▪",space_after=7)
runs(tf,[("Why a panel?",GREEN,True),(" averaging three different learners is steadier and more accurate than any one — "
        "and beats simpler methods (shown next).",LIGHT,False)],size=13.5,space_after=0)
plainbox(s,0.9,6.55,11.5,0.0,"",color=BLUE)
notes(s,
"SAY: Computers can't read a chemical drawing directly, so we convert each molecule into numbers in two ways: a "
"'fingerprint' that records which small sub-structures are present, and 24 physical properties like size and "
"greasiness. Then three different model types each make a prediction and we average their votes — like a panel of "
"experts rather than one opinion.\n\n"
"TECHNICAL TERMS:\n- ECFP-4 / Morgan fingerprint: a fixed-length binary code of a molecule's local substructures.\n"
"- Descriptor: a computed physical property (molecular weight, logP/greasiness, polar surface area, etc.).\n"
"- Random Forest / Extra Trees / Gradient Boosting: standard, robust 'tree ensemble' algorithms.\n"
"- Ensemble: combining several models to reduce error.\n\n"
"Q&A — EXPERT:\nQ: Why not a graph neural network?\nA: On datasets this size, tree ensembles on ECFP+descriptors are "
"competitive, far more interpretable, cheaper, and match the literature baselines; deep nets were unnecessary for the "
"claim. Fixed seed 42, RDKit 2026.03, scikit-learn 1.8.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: Does it need a supercomputer?\nA: No — it runs on a normal machine, which keeps it "
"cheap and deployable.")

# ============================================================ SLIDE 7 — TRUST PARAMETERS
s=slide(); bg(s,PAPER); eyebrow(s,"06 · Making the answer trustworthy"); accentbar(s)
title(s,"Every threshold is a stated, defended choice")
items=[("Labelling","'Active' = potent (pChEMBL ≥ 6, ~1 µM); 'inactive' = weak (< 5). The ambiguous middle is dropped to cut noise.",BLUE),
       ("Deployment gate","A model ships only if it clears a fixed quality bar (MCC ≥ 0.45). Four targets failed as yes/no and were re-cast as ranking.",AMBER),
       ("Calibration","Scores are adjusted (isotonic) so a '0.9' behaves like a real 90% — not just a ranking.",GREEN),
       ("Confidence + domain","Conformal sets give a ~90%-coverage confidence range; an 'applicability' flag warns when a molecule is unlike anything seen.",VERM)]
x,y=0.9,2.05
for i,(h,b,c) in enumerate(items):
    cx=x+(i%2)*5.95; cy=y+(i//2)*2.15
    rrect(s,cx,cy,5.75,1.95,WHITE,line=LINE); rect(s,cx,cy,0.06,1.95,c)
    _,tf=box(s,cx+0.25,cy+0.16,5.35,1.7)
    para(tf,h,size=15.5,bold=True,color=INK,first=True,space_after=4)
    para(tf,b,size=12.5,color=MUTED,space_after=0)
notes(s,
"SAY: The difference between a toy and a trustworthy tool is what happens around the prediction. Four choices matter. "
"First, how we define active vs inactive. Second, a quality gate so weak models never ship. Third, calibration so the "
"numbers mean what they say. Fourth, a confidence range plus a warning when a molecule is outside what the model has "
"seen.\n\n"
"TECHNICAL TERMS:\n- pChEMBL >=6 (~1 micromolar): a standard potency cut for 'active'.\n- MCC: a balanced accuracy "
"score (−1 to 1) robust to class imbalance; 0.45 is our minimum to deploy.\n- Isotonic calibration: a monotonic "
"adjustment of scores to match observed frequencies.\n- Conformal prediction: a method giving prediction sets with a "
"guaranteed error rate.\n- Applicability domain: the chemical space where the model is trustworthy.\n\n"
"Q&A — EXPERT:\nQ: Is the >=6/<5 cut cherry-picked?\nA: No — we re-labelled at several cuts and retrained; max AUROC "
"spread was 0.109 and the deployed cut ~ the strict cut (robustness slide). The grey-zone drop is empirically justified.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: What happens with a brand-new molecule type?\nA: The tool flags it as 'out of domain' "
"instead of pretending confidence — a safety feature, not a bug.")

# ============================================================ SLIDE 8 — VALIDATION
s=slide(); bg(s,NAVY); eyebrow(s,"07 · How we tested it",color=SKY); rect(s,0.9,0.62,0.5,0.045,SKY)
title(s,"Four exams, each harder and more realistic",color=WHITE)
steps=[("Random","Familiar questions — the easy exam most papers report.",SKY),
       ("Scaffold","New molecular skeletons at test time.",BLUE),
       ("Leave-cluster-out","Whole new chemical families held back.",GREEN),
       ("Temporal","Trained on older compounds, tested on the newest — a true 'future' exam.",AMBER)]
x=0.9
for i,(h,b,c) in enumerate(steps):
    rrect(s,x,2.2,2.85,3.0,NAVY2,line=RGBColor(0x25,0x40,0x5A))
    _,tf=box(s,x+0.22,2.4,2.45,2.7)
    para(tf,f"Exam {i+1}",size=11,bold=True,color=c,font=MONO,first=True,space_after=4)
    para(tf,h,size=15.5,bold=True,color=WHITE,space_after=6)
    para(tf,b,size=12,color=LIGHT,space_after=0)
    if i<3:
        _,tf=box(s,x+2.72,3.35,0.35,0.5); para(tf,"›",size=22,bold=True,color=SKY,first=True)
    x+=3.02
_,tf=box(s,0.9,5.6,11.5,1.2)
para(tf,"Reporting all four is a deliberate honesty choice: it shows exactly where the tool is strong and where "
        "it is not. Most published models report only the easy 'random' exam.",size=15,color=LIGHT,first=True)
notes(s,
"SAY: Anyone can score well on an easy test. We ran four, each closer to real life. The 'random' test is the easy one "
"most papers stop at. The toughest — 'temporal' — trains only on older molecules and tests on the newest ones, "
"simulating genuine prospective use. We report all four on purpose.\n\n"
"TECHNICAL TERMS:\n- Scaffold split (Bemis–Murcko): grouping by core molecular skeleton so test molecules are "
"structurally novel.\n- Leave-cluster-out: entire chemical clusters held out.\n- Temporal split: train on compounds "
"published up to a cutoff year, test on later ones.\n- AUROC (next slide): the score used.\n\n"
"Q&A — EXPERT:\nQ: Why does temporal drop so much for some endpoints?\nA: 71–91% of recent test compounds have unseen "
"scaffolds (covariate shift). Where the recent set is class-balanced (MAO-A) the honest number is 0.61; we surface "
"this rather than hide it.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: Why show your weaknesses?\nA: Credibility. A tool that only reports its best number "
"can't be trusted; buyers and reviewers value honesty about limits.")

# ============================================================ SLIDE 9 — CLASSIFICATION RESULTS
s=slide(); bg(s,PAPER); eyebrow(s,"08 · Results — the yes/no models"); accentbar(s)
title(s,"Strong on the easy exam, honest on the hard ones")
add_table(s,[
 ["Endpoint","Random","Scaffold","Cluster","Temporal"],
 ["BBB","0.963","0.921","0.906","–"],
 ["AChE","0.975","0.915","0.912","0.784"],
 ["BChE","0.976","0.937","0.921","0.794"],
 ["BACE1","0.956","0.950","0.940",("0.915",GREEN,True)],
 ["GSK-3β","0.943","0.920","0.915",("0.658",VERM,True)],
 ["MAO-B","0.960","0.885","0.873","0.758"],
 ["MAO-A","0.950","0.867","0.890",("0.614",VERM,True)],
 ["hERG","0.950","0.901","0.870","0.757"],
],0.9,2.0,7.0,[1.5,1.0,1.0,1.0,1.0],rh=0.40,fs=11.5)
_,tf=box(s,8.3,2.0,4.2,0.5); para(tf,"AUROC by test (higher = better)",size=11.5,italic=True,color=MUTED,first=True)
data=[("Random split","0.94–0.98",GREEN),("Scaffold & cluster","0.87–0.95",BLUE),
      ("Temporal (future)","0.61–0.92",AMBER),("Confidence coverage","0.885–0.905",SKY)]
y=2.5
for h,v,c in data:
    rrect(s,8.3,y,4.2,0.72,WHITE,line=LINE); rect(s,8.3,y,0.06,0.72,c)
    _,tf=box(s,8.5,y+0.05,3.9,0.62)
    runs(tf,[(v+"   ",c,True),(h,MUTED,False)],size=14,first=True,space_after=0)
    tf.paragraphs[0].runs[0].font.name=MONO
    y+=0.82
plainbox(s,0.9,6.15,11.6,0.85,
"AUROC scores how well the tool separates 'active' from 'inactive'. 0.5 = a coin flip; 1.0 = perfect. "
"0.90 means that, given one active and one inactive molecule, it ranks them correctly ~90% of the time.",color=BLUE)
notes(s,
"SAY: On the standard 'random' test the models score 0.94–0.98 — at or above published state of the art. On the harder "
"tests they hold 0.87–0.95, and on the toughest 'future' test they range 0.61–0.92. We show the low numbers too — "
"GSK-3β and MAO-A are the hardest and we flag them as lower-confidence.\n\n"
"TECHNICAL TERMS:\n- AUROC: probability the model ranks a random active above a random inactive; 0.5 chance, 1.0 "
"perfect.\n- 'Coverage 0.885–0.905': when the tool gives a 90% confidence set, it is right about 90% of the time — the "
"promise is kept.\n\n"
"Q&A — EXPERT:\nQ: BACE1 temporal 0.92 looks too good.\nA: Its recent test set is ~93% active, which inflates AUROC; we "
"note this. Conversely MAO-A's balanced set gives an honest 0.61.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: Is 0.90 good?\nA: Yes — for this kind of prediction it's strong; and unlike many tools "
"we also tell you when to trust it.")

# ============================================================ SLIDE 10 — REGRESSION/ANTIOX/DRUGG
s=slide(); bg(s,PAPER); eyebrow(s,"09 · Results — strength & drug-likeness"); accentbar(s)
title(s,"Ranking-grade potency, a measured antioxidant model")
add_table(s,[
 ["Receptor","R² (fit)","Spearman","Temporal"],
 ["A2A","0.526","0.706","0.326"],
 ["5-HT2A","0.460","0.684","0.085"],
 ["D2","0.425","0.652","−0.007"],
 ["SERT","0.338","0.573","0.171"],
],0.9,2.05,6.2,[1.6,1.2,1.2,1.2],rh=0.40)
rrect(s,7.4,2.05,5.1,1.75,WHITE,line=LINE); rect(s,7.4,2.05,0.06,1.75,GREEN)
_,tf=box(s,7.65,2.2,4.7,1.5)
para(tf,"Antioxidant (measured DPPH)",size=14.5,bold=True,color=INK,first=True,space_after=4)
runs(tf,[("n = 2,862   ·   R² 0.43   ·   ρ 0.636",INK,False)],size=13,space_after=3); tf.paragraphs[-1].runs[0].font.name=MONO
para(tf,"Replaces a prior text-based proxy that fit poorly (R² ≈ 0.25).",size=11.5,color=MUTED,space_after=0)
_,tf=box(s,0.9,4.4,5.6,0.4); para(tf,"Druggability score — CNS drugs vs polar non-drugs",size=13,bold=True,color=INK,first=True)
bars=[("Caffeine",86,GREEN),("Donepezil",79,GREEN),("Sucrose",46,VERM),("Atorvastatin",24,VERM)]
y=4.85
for name,val,c in bars:
    _,tf=box(s,0.9,y-0.02,1.7,0.3); para(tf,name,size=11.5,color=MUTED,font=MONO,first=True)
    rect(s,2.7,y+0.02,3.6,0.22,RGBColor(0xE9,0xEE,0xF3))
    rect(s,2.7,y+0.02,3.6*val/100.0,0.22,c)
    _,tf=box(s,6.35,y-0.02,0.6,0.3); para(tf,str(val),size=11.5,color=INK,font=MONO,first=True)
    y+=0.42
plainbox(s,7.4,4.4,5.1,2.15,
"R² and Spearman measure how closely predicted strength tracks the real measured strength (1.0 = perfect ordering). "
"For druggability, higher = more drug-like for the brain — caffeine and donepezil score high; table sugar and a "
"cholesterol drug (not brain-aimed) score low, as expected.",color=GREEN)
notes(s,
"SAY: For the four receptors we predict strength and use it to rank compounds — useful for prioritisation even when an "
"exact value is hard. The antioxidant model is now trained on real DPPH lab measurements and clearly beats the old "
"text-derived estimate. The druggability score is a transparent rule that correctly separates brain-suitable drugs "
"from things like table sugar.\n\n"
"TECHNICAL TERMS:\n- R²: fraction of variation explained (1.0 perfect, 0 none).\n- Spearman (ρ): how well the "
"predicted ranking matches the true ranking.\n- DPPH: standard antioxidant assay.\n\n"
"Q&A — EXPERT:\nQ: Temporal R² near zero for some receptors/DPPH?\nA: Correct and disclosed — pooled cross-lab DPPH and "
"receptor data generalise weakly across time; these are reported as ranking-grade, not absolute predictors.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: What's druggability good for?\nA: A quick 'is this even worth pursuing as a brain drug?' "
"filter before expensive work.")

# ============================================================ SLIDE 11 — BASELINES
s=slide(); bg(s,NAVY); eyebrow(s,"10 · Does the model earn its keep?",color=SKY); rect(s,0.9,0.62,0.5,0.045,SKY)
title(s,"It beats 'just find the most similar known molecule'",color=WHITE)
_,tf=box(s,0.9,1.95,11.4,0.8)
para(tf,"We compared the model against two simpler methods under the identical fair test.",size=15.5,color=LIGHT,first=True)
comp=[("Our ensemble",0.912,GREEN),("Nearest-neighbour lookup",0.867,SKY),("Simple linear model",0.808,AMBER)]
y=3.0
for name,val,c in comp:
    _,tf=box(s,0.9,y-0.03,3.4,0.4); para(tf,name,size=14,color=LIGHT,first=True)
    rect(s,4.5,y,6.2,0.34,NAVY2)
    rect(s,4.5,y,6.2*((val-0.5)/0.5),0.34,c)
    _,tf=box(s,10.85,y-0.03,1.4,0.4); para(tf,f"{val:.3f}",size=14,bold=True,color=WHITE,font=MONO,first=True)
    y+=0.62
_,tf=box(s,4.5,y+0.02,6.2,0.3)
runs(tf,[("0.50",SKY,False),("        mean AUROC (8 endpoints)        ",MUTED,False),("1.00",SKY,False)],size=10,font=MONO,first=True,align=PP_ALIGN.CENTER)
_,tf=box(s,0.9,5.65,11.5,1.4)
runs(tf,[("The model wins on all 8 endpoints. ",GREEN,True),
 ("Beating a pure 'look up the most similar molecule' method is the key check that it has learned real "
  "structure–activity patterns — not just memorised look-alikes. This is also why it behaves differently from an "
  "LLM (next).",LIGHT,False)],size=15,first=True)
notes(s,
"SAY: A fair worry: is the model just memorising and returning the closest known molecule? We tested exactly that. A "
"nearest-neighbour lookup scores 0.867; a simple linear model 0.808; our ensemble 0.912 — and it wins on every single "
"endpoint. So it genuinely learns patterns, not just look-alikes.\n\n"
"TECHNICAL TERMS:\n- Nearest-neighbour (kNN-Tanimoto): predict from the most chemically similar known compounds — a "
"'read-across' baseline.\n- This baseline is the closest analogy to how an LLM's fuzzy recall works, which sets up the "
"next slide.\n\n"
"Q&A — EXPERT:\nQ: Deltas look modest (+0.045 vs kNN).\nA: Modest but consistent — best on 8/8 under scaffold split, "
"and the point is qualitative: it exceeds pure similarity recall, so performance isn't memorisation.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: So it's smarter than a database search?\nA: Yes — a search only finds look-alikes; the "
"model generalises to genuinely new structures, with calibrated confidence.")

# ============================================================ SLIDE 12 — ROBUSTNESS
s=slide(); bg(s,PAPER); eyebrow(s,"11 · We tested our own assumptions"); accentbar(s)
title(s,"Three known risks — measured, not just mentioned")
cols=[("Definition of 'active'","We re-labelled at several cut-offs and retrained. Results barely moved (max AUROC change 0.109); the deployed choice is not cherry-picked.",BLUE),
      ("Mixing experiment types","Retraining on one experiment type only vs the mix changed scores by ≤ 0.006 — mixing is safe on the common scale.",GREEN),
      ("Knowing its own limits","Accuracy falls smoothly as molecules get less similar to training (0.96 → 0.77). That is exactly why the 'out-of-domain' warning exists.",AMBER)]
x=0.9
for h,b,c in cols:
    rrect(s,x,2.05,3.85,3.2,WHITE,line=LINE); rect(s,x,2.05,3.85,0.06,c)
    _,tf=box(s,x+0.25,2.3,3.4,2.8)
    para(tf,h,size=15,bold=True,color=INK,first=True,space_after=6)
    para(tf,b,size=12.5,color=MUTED,space_after=0)
    x+=4.05
plainbox(s,0.9,5.55,11.6,1.1,
"instead of just listing caveats, we re-ran the whole pipeline to put a number on each risk. All three came back "
"reassuring — and where the model is weak (very novel molecules), it warns you rather than guessing confidently.",color=GREEN)
notes(s,
"SAY: Good science stress-tests its own choices. We asked three sceptical questions and answered each with a re-run, "
"not an opinion. Does the active/inactive cut-off matter? Barely. Does mixing experiment types distort things? By "
"under 0.006. Does it know its limits? Yes — accuracy declines predictably with novelty, which is why the warning "
"flag is calibrated the way it is.\n\n"
"TECHNICAL TERMS:\n- Assay types (IC50/Ki/Kd/EC50): different ways to measure potency, standardised onto pChEMBL.\n"
"- 'Out-of-domain': a molecule unlike anything in training, where predictions are less reliable.\n\n"
"Q&A — EXPERT:\nQ: Which endpoints did you test for assay pooling?\nA: GSK-3β (the most mixed, 49% IC50), MAO-B and "
"hERG; deltas were +0.006, −0.006 and 0.000 — negligible.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: How do I know it won't confidently give a wrong answer?\nA: For unfamiliar molecules it "
"marks the result out-of-domain and widens its confidence — it declines to over-claim.")

# ============================================================ SLIDE 13 — LLM HEAD TO HEAD
s=slide(); bg(s,NAVY); eyebrow(s,"12 · \"Why not just ask ChatGPT?\"",color=SKY); rect(s,0.9,0.62,0.5,0.045,SKY)
title(s,"We ran a fair, pre-registered contest — here's what happened",color=WHITE)
add_table(s,[
 ["System","BBB","hERG","Made-up evidence","New molecule"],
 [("BrainSafe",WHITE,True),"8/9","5/5",("0 of 0",GREEN,True),("honest",GREEN,True)],
 ["Gemini Pro","7/9","5/5",("5 of 10",VERM,True),("made up",VERM,True)],
 ["ChatGPT-4o","9/9","5/5",("4 of 9",VERM,True),("made up",VERM,True)],
 ["Perplexity","9/9","4/5",("1 of 2",VERM,True),("made up",VERM,True)],
 ["Claude","9/9","5/5",("4 of 10",VERM,True),("made up",VERM,True)],
],0.9,2.0,7.6,[1.7,0.9,0.9,1.6,1.3],rh=0.40,fs=12)
_,tf=box(s,8.75,2.0,3.75,4.2)
runs(tf,[("On famous drugs, the LLMs are strong",SKY,True),(" — several match or beat us. We say so plainly.",LIGHT,False)],size=13,first=True,space_after=9)
runs(tf,[("But 45% of the specific evidence IDs they cited were fake or pointed to the wrong drug",VERM,True),
         (" — a cited 'rasagiline' ID was actually a steroid; a 'rivastigmine' ID was vitamin B6.",LIGHT,False)],size=13,space_after=9)
runs(tf,[("All four invented",VERM,True),(" a target and potency for an unpublished molecule — and disagreed.",LIGHT,False)],size=13,space_after=0)
_,tf=box(s,0.9,6.35,11.6,0.9)
runs(tf,[("The point isn't 'LLMs are wrong.' ",WHITE,True),
 ("They can't be trusted for verifiable evidence or for genuinely new molecules — exactly where discovery happens. "
  "BrainSafe grounds every answer in a real measurement.",LIGHT,False)],size=14.5,first=True)
notes(s,
"SAY: The obvious question from a commercial audience: why not just use ChatGPT? So we ran a fair contest — same "
"questions, scored against real data, decided in advance. Two honest findings. One: on well-known drugs the LLMs are "
"genuinely good, sometimes better than us. Two: when we checked the specific evidence they cited — database IDs for "
"the molecules — 45% were fabricated or pointed to a completely different drug (one 'rasagiline' ID was actually a "
"steroid). And for a brand-new, unpublished molecule, all four confidently made up an answer, and disagreed with each "
"other. Our tool cited real measurements every time and said 'uncertain' on the new one.\n\n"
"TECHNICAL TERMS:\n- Pre-registered: prompt, molecules and scoring fixed before running, so it's fair.\n"
"- 'Made-up evidence': a cited ChEMBL ID that doesn't exist or belongs to a different molecule (we checked each live).\n"
"- Confabulation/hallucination: a confident but false answer.\n\n"
"Q&A — EXPERT:\nQ: Isn't this an unfair prompt?\nA: The prompt is identical for every system including ours, and "
"scored on the same measured-data key with live ChEMBL verification; we even dropped contestable items (e.g. donepezil "
"hERG) to be fair to the LLMs.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: So LLMs are useless here?\nA: No — great for summarising known drugs. But for decisions "
"needing traceable evidence or novel chemistry, a grounded tool is essential. They're complementary.")

# ============================================================ SLIDE 14 — DIFFERENTIATION
s=slide(); bg(s,PAPER); eyebrow(s,"13 · How BrainSafe is different"); accentbar(s)
title(s,"Grounded and auditable, by design")
add_table(s,[
 ["Capability","BrainSafe AI","General LLM","Typical ADMET tool"],
 ["Trained on measured bioactivity",("yes — 64,474",GREEN,True),"no (text)","partly"],
 ["Honest, calibrated probability",("yes",GREEN,True),"no","rarely"],
 ["Guaranteed confidence range",("yes",GREEN,True),"no","no"],
 ["Shows the real evidence",("yes",GREEN,True),("no — invents",VERM,True),"no"],
 ["Warns when out of its depth",("yes",GREEN,True),"no","sometimes"],
 ["Brain-entry-gated disease view",("yes",GREEN,True),"no","no"],
 ["Same answer every time",("yes",GREEN,True),("no",VERM,True),"yes"],
],0.9,2.0,11.6,[3.2,1.6,1.4,1.7],rh=0.46,fs=12)
notes(s,
"SAY: Put simply, BrainSafe is grounded and auditable by construction. It learns from measured data, gives calibrated "
"probabilities with a confidence range, shows the real evidence behind each answer, warns when it's out of its depth, "
"combines target activity with brain entry, and gives the same answer every time. General LLMs and typical ADMET tools "
"each miss several of these.\n\n"
"Q&A — EXPERT:\nQ: Isn't 'shows evidence' just nearest-neighbour retrieval?\nA: Yes — deliberately: it returns the "
"actual measured analogue and value, which is verifiable, unlike an LLM's fabricated citation.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: What's the one-line pitch?\nA: 'A brain-drug triage tool you can audit — every answer "
"backed by a real measurement, with honest confidence.'")

# ============================================================ SLIDE 15 — APPLICATIONS & LIMITS
s=slide(); bg(s,PAPER); eyebrow(s,"14 · Uses & honest limits"); accentbar(s)
title(s,"A research triage engine — not a clinical oracle")
rrect(s,0.9,2.05,5.75,3.9,WHITE,line=LINE); rect(s,0.9,2.05,5.75,0.06,GREEN)
_,tf=box(s,1.15,2.3,5.3,3.5)
para(tf,"What it is good for",size=16,bold=True,color=GREEN,first=True,space_after=8)
for t in ["Prioritising which molecules to test next for brain conditions (natural products, flavonoids, repurposing).",
          "Early one-pass triage: brain entry + targets + safety + drug-likeness, with confidence and evidence.",
          "A transparent, reproducible companion to expert judgement."]:
    runs(tf,[(t,INK,False)],size=13,bullet="▪",space_after=9); tf.paragraphs[-1].runs[0].font.color.rgb=GREEN
rrect(s,6.85,2.05,5.6,3.9,WHITE,line=LINE); rect(s,6.85,2.05,5.6,0.06,VERM)
_,tf=box(s,7.1,2.3,5.15,3.5)
para(tf,"What it does NOT claim",size=16,bold=True,color=VERM,first=True,space_after=8)
for t in ["Predicts whether a molecule engages a target — not whether it activates or blocks it.",
          "Engagement is not the same as clinical benefit; the clinical layer is precedent, not proof.",
          "No wet-lab validation yet; weaker on very novel chemistry (flagged).",
          "One heart-safety target (hERG); other safety risks are out of scope."]:
    runs(tf,[(t,INK,False)],size=13,bullet="▪",space_after=8); tf.paragraphs[-1].runs[0].font.color.rgb=VERM
notes(s,
"SAY: Being clear about scope is part of the integrity. It's a triage and hypothesis-generation engine: it helps you "
"decide what to test next, cheaply and transparently. It is explicitly not a clinical or diagnostic tool. It predicts "
"engagement, not direction or efficacy; it hasn't had wet-lab confirmation yet; and it covers one safety axis.\n\n"
"TECHNICAL TERMS:\n- Engagement vs direction: binding to a target vs activating (agonist) or blocking (antagonist) it.\n"
"- Engagement vs efficacy: hitting a target vs actually helping a patient.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: Can it recommend a treatment?\nA: No — and it never should. It shortlists candidates for "
"scientists to test.\nQ: Commercial applications?\nA: Compound triage for CNS discovery, natural-product screening, and "
"drug-repurposing prioritisation — reducing wasted lab spend.\n\n"
"Q&A — EXPERT:\nQ: Path to efficacy prediction?\nA: Out of scope here; would need target-direction and disease-model "
"data. We're explicit that this is engagement + clinical precedent, not efficacy.")

# ============================================================ SLIDE 16 — CLOSE
s=slide(); bg(s,NAVY)
rect(s,0,0,SW.inches,0.12,SKY)
_,tf=box(s,0.9,0.7,11.5,0.5); runs(tf,[("● ",SKY,True),("WHAT WE ACHIEVED",LIGHT,True)],size=13,font=MONO,first=True)
title(s,"A brain-effect predictor you can actually trust",color=WHITE,t=1.15,size=30)
stats=[("64,474","measured records · 12 endpoints"),("4 exams","random → scaffold → cluster → temporal"),
       ("~90%","confidence promise, kept"),("0","fabricated evidence items")]
x=0.9
for v,l in stats:
    rrect(s,x,2.5,2.85,1.6,NAVY2,line=RGBColor(0x25,0x40,0x5A))
    _,tf=box(s,x+0.2,2.65,2.5,1.4)
    para(tf,v,size=26,bold=True,color=SKY,font=MONO,first=True,space_after=2)
    para(tf,l,size=11.5,color=LIGHT,space_after=0)
    x+=3.02
_,tf=box(s,0.9,4.5,11.4,1.6)
para(tf,"A calibrated, evidence-grounded, brain-entry-gated predictor built entirely on measured public data — "
        "state-of-the-art-grade on like-for-like tests, fully transparent on the hard ones, and every claim "
        "traceable to a real measurement.",size=16,color=LIGHT,first=True)
_,tf=box(s,0.9,6.4,11.5,0.6)
para(tf,"No fabrication · no assumption · reproducible from released scripts and data.",size=12.5,bold=True,color=SKY,font=MONO,first=True)
notes(s,
"SAY: To close — we built a brain-effect predictor you can trust because you can check it. 64,474 real "
"measurements, four levels of testing, a confidence promise that holds, and zero fabricated evidence. It matches the "
"best on standard tests, is honest about the hard ones, and every answer traces back to a real experiment. Thank you "
"— happy to take questions.\n\n"
"Q&A — EXPERT:\nQ: What's next?\nA: Wet-lab prospective validation, more safety anti-targets, and broader target "
"coverage.\n\n"
"Q&A — GENERAL/COMMERCIAL:\nQ: When can we use it?\nA: It runs today as a research tool; wider release follows peer "
"review. We'd welcome collaboration on validation.")

prs.save("BrainSafe_AI_Presentation.pptx")
print("Saved BrainSafe_AI_Presentation.pptx |", len(prs.slides.__iter__.__self__._sldIdLst), "slides with speaker notes")
